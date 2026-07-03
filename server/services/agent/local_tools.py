"""本地文件操作工具（桌面全本地版专属，DESKTOP_LOCAL=1 才注册）。

把"像 Claude Code 那样在本机读/写/改文件"的通用能力给 Agent——但配台球老板用得起的护栏：
- **范围锁**：只动「内容库」(用户数据目录下一个文件夹) + 库内文件。绝不漫游全盘、不碰系统文件。
- **改前备份**：写/改前自动把原件复制到 .backups/，可回滚。
- **审批闸**：write/edit 类 requires_approval=True，循环里不直接执行——先把"要怎么改"弹给老板(人看得懂的改动)，确认后才落盘。
- **不暴露裸 shell**：只给文件操作（改动是人能审的 diff），命令类另包成具体安全动作。

⚠️ 云端 web 版（PostgreSQL，多租户）绝不注册这些——文件操作只在用户自己机器上的本地后端有意义。
"""
import asyncio
import hashlib
import json
import logging
import os
import re
import shlex
import shutil
import signal
import subprocess
from datetime import datetime
from pathlib import Path

from core.timezone import business_now
from services.agent.registry import Tool, default_registry

logger = logging.getLogger(__name__)


def _library_root() -> Path:
    """内容库根：Agent 只能在这里面(+库内)动手。默认在用户数据目录下。"""
    root = Path(os.environ.get("DESKTOP_LIBRARY_DIR") or (Path.home() / ".billiards-desktop" / "library"))
    root.mkdir(parents=True, exist_ok=True)
    return root


def _allowed_paths(ctx) -> list[Path]:
    """用户当场选定、显式授权的文件/目录（来自 OS 文件选择器）。解析为绝对 Path。"""
    raw = getattr(ctx, "allowed_paths", None) or []
    out: list[Path] = []
    for s in raw:
        try:
            out.append(Path(s).resolve())
        except (OSError, ValueError):
            continue
    return out


def _resolve(rel_or_abs: str, ctx=None) -> Path:
    """把传入路径解析进沙箱并校验不越界。沙箱 = 内容库 + 工作目录 + 用户当场选定的文件/目录。
    返回绝对 Path；越界抛 ValueError。相对路径落工作目录(没设则落内容库)。
    ctx.full_disk_access=True 时不限范围(桌面应用默认 True；沙箱模式才传 False·对标 CC 全盘可达)。"""
    lib = _library_root().resolve()
    wd = None
    raw_wd = (getattr(ctx, "working_dir", None) or "").strip() or None  # 与 Task1 一致:纯空白当没设
    if raw_wd:
        try:
            wd = Path(raw_wd).resolve()
        except (OSError, ValueError):
            wd = None
    base = wd or lib
    p = Path(rel_or_abs)
    path = (p if p.is_absolute() else base / p).resolve()
    if getattr(ctx, "full_disk_access", False):
        return path
    # 内容库内 / 工作目录内 → 放行
    for rootp in (lib, wd):
        if rootp and (path == rootp or rootp in path.parents):
            return path
    # 用户经文件选择器当场选定的文件/目录(或其子文件)→ 放行(显式授权)
    for a in _allowed_paths(ctx):
        if path == a or a in path.parents:
            return path
    # Agent 自己落盘的超大工具结果目录 → 放行
    try:
        from services.agent.tool_result_store import results_root
        tr_root = results_root().resolve()
        if path == tr_root or tr_root in path.parents:
            return path
    except Exception:
        pass
    raise ValueError(f"越界：只能操作内容库、工作目录或你当场选定的文件，拒绝 {rel_or_abs}")


def path_in_workspace(raw_path: str, ctx) -> bool:
    """目标路径是否在【工作区】内 = 内容库 ∪ 工作目录 ∪ 用户选定文件/目录 ∪ tool-results。
    相对路径按 (ctx.working_dir or 内容库) 解析。用于 auto_files 判"改这文件要不要免确认"。
    不抛、故障安全(判不了→False=更安全的弹卡)。与 full_disk_access 无关(它管免不免确认、不管能不能碰)。"""
    if not raw_path:
        return False
    try:
        lib = _library_root().resolve()
        wd = None
        raw_wd = (getattr(ctx, "working_dir", None) or "").strip() or None  # 与 _resolve 一致：纯空白当没设
        if raw_wd:
            try:
                wd = Path(raw_wd).resolve()
            except (OSError, ValueError):  # 与 _resolve 一致：坏路径优雅降级落内容库，而非整体返 False
                wd = None
        base = wd or lib
        p = Path(raw_path)
        path = (p if p.is_absolute() else base / p).resolve()
        roots = [lib] + ([wd] if wd else []) + _allowed_paths(ctx)
        try:
            from services.agent.tool_result_store import results_root
            roots.append(results_root().resolve())
        except Exception:
            pass
        for r in roots:
            if r and (path == r or r in path.parents):
                return True
        return False
    except Exception:
        return False


def is_path_allowed(raw_path: str, ctx) -> bool:
    """审批闸 2.0 · ① 薄封装：给 loop 用的"这条路径 `_resolve` 时会不会越界抛错"判定（不抛错版）。

    语义必须与 `_resolve` 完全一致——先看 `full_disk_access`（开了就不限范围，任意路径放行，
    与 `_resolve` 的门控行为对齐），没开才落到 `path_in_workspace`（内容库/工作目录/用户选定/
    tool-results 白名单）。特意薄封装成一个函数，让 loop.py 判"文件类工具目标是否越界"时直接调用，
    不在 loop 里重抄一遍沙箱规则（沙箱规则的唯一真相源永远在这个模块）。"""
    if getattr(ctx, "full_disk_access", False):
        return True
    return path_in_workspace(raw_path, ctx)


def _path_hash(path: Path) -> str:
    """路径 hash 前缀（8 位），防不同目录同名文件的备份撞名。"""
    return hashlib.md5(str(path.resolve()).encode()).hexdigest()[:8]


def _backup_meta_path(backup: Path) -> Path:
    return backup.with_suffix(backup.suffix + ".json")


def _read_backup_meta(backup: Path) -> dict:
    try:
        meta = json.loads(_backup_meta_path(backup).read_text(encoding="utf-8"))
        return meta if isinstance(meta, dict) else {}
    except Exception:
        return {}


def _same_resolved_path(left: str | Path, right: str | Path) -> bool:
    try:
        return Path(left).expanduser().resolve() == Path(right).expanduser().resolve()
    except (OSError, ValueError):
        return False


def _backup(path: Path) -> str | None:
    """改/写前备份原件，返回备份路径（原件不存在则 None）。
    备份名带路径 hash 前缀，不同目录的同名文件不会互相覆盖。"""
    if not path.exists():
        return None
    bdir = _library_root() / ".backups"
    bdir.mkdir(exist_ok=True)
    stamp = business_now().strftime("%Y%m%d-%H%M%S-%f")
    prefix = _path_hash(path)
    dest = bdir / f"{prefix}_{path.stem}.{stamp}{path.suffix}.bak"
    shutil.copy2(path, dest)
    try:
        _backup_meta_path(dest).write_text(json.dumps({
            "original_path": str(path.resolve()),
            "backup_path": str(dest.resolve()),
            "created_at": business_now().isoformat(),
            "size": path.stat().st_size,
        }, ensure_ascii=False, indent=2), encoding="utf-8")
    except Exception:
        logger.debug("failed to write backup metadata for %s", dest, exc_info=True)
    return str(dest)


def _backup_candidates(path: Path) -> list[Path]:
    bdir = _library_root() / ".backups"
    if not bdir.is_dir():
        return []
    prefix = _path_hash(path)
    cands = sorted(bdir.glob(f"{prefix}_{path.stem}.*{path.suffix}.bak"))
    if not cands:
        cands = sorted(bdir.glob(f"{path.stem}.*{path.suffix}.bak"))
    return [p for p in cands if p.is_file()]


def list_file_backups(limit: int = 20) -> list[dict]:
    """列出最近文件备份，供前端轻量找回“AI 改过/删过的文件”。"""
    bdir = _library_root() / ".backups"
    if not bdir.is_dir():
        return []
    rows: list[dict] = []
    for backup in bdir.glob("*.bak"):
        if not backup.is_file():
            continue
        meta = _read_backup_meta(backup)
        original = meta.get("original_path") or ""
        rows.append({
            "backup_path": str(backup),
            "path": original,
            "name": Path(original).name if original else backup.name.replace(".bak", ""),
            "created_at": meta.get("created_at") or datetime.fromtimestamp(backup.stat().st_mtime).isoformat(),
            "size": meta.get("size") or backup.stat().st_size,
            "exists": Path(original).exists() if original else False,
        })
    rows.sort(key=lambda x: str(x.get("created_at") or ""), reverse=True)
    return rows[:max(1, min(int(limit or 20), 80))]


def get_file_backup_diff(raw_path: str, backup_path: str | None = None) -> dict:
    """B.2：给一个被 AI 改过的本机文件，返回 {ok, path, old, new}——old=最近一次备份、new=当前内容，
    供前端右侧渲染"改前/改后"对比让老板确认。只读、故障安全。找不到备份→old=""(当新建处理)。"""
    try:
        p = Path(raw_path).expanduser()
        if p.exists() and not p.is_file():
            return {"ok": False, "error": "这不是一个文件"}
        if not p.exists() and backup_path is None:
            return {"ok": False, "error": "文件不在了"}
        if p.exists() and p.stat().st_size > 2 * 1024 * 1024:
            return {"ok": False, "error": "文件太大，不便逐字对比"}
        new = p.read_text(encoding="utf-8", errors="replace") if p.is_file() else ""
        old = ""
        backup = Path(backup_path).expanduser() if backup_path else None
        if backup is not None:
            bdir = (_library_root() / ".backups").resolve()
            try:
                backup = backup.resolve()
            except (OSError, ValueError):
                return {"ok": False, "error": "备份路径不对"}
            if backup.parent != bdir or backup.suffix != ".bak":
                return {"ok": False, "error": "只能查看自动备份里的文件"}
            meta = _read_backup_meta(backup)
            original = meta.get("original_path")
            if original and not _same_resolved_path(original, p):
                return {"ok": False, "error": "备份和文件不匹配"}
            cands = [backup] if backup.is_file() else []
        else:
            cands = _backup_candidates(p)
        if cands:
            try:
                old = cands[-1].read_text(encoding="utf-8", errors="replace")
                backup = cands[-1]
            except Exception:
                old = ""
        return {"ok": True, "path": str(p), "backup_path": str(backup) if backup else None, "old": old, "new": new}
    except Exception as e:
        return {"ok": False, "error": str(e)[:120]}


def restore_file_backup(raw_path: str, backup_path: str | None = None) -> dict:
    """把某个文件恢复到最近备份。恢复前也备份当前版本，避免“回退”本身不可逆。"""
    try:
        p = Path(raw_path).expanduser()
        backup = Path(backup_path).expanduser() if backup_path else None
        if backup is None:
            cands = _backup_candidates(p)
            if not cands:
                return {"ok": False, "error": "没有找到这份文件的备份"}
            backup = cands[-1]
        bdir = (_library_root() / ".backups").resolve()
        try:
            backup = backup.resolve()
        except (OSError, ValueError):
            return {"ok": False, "error": "备份路径不对"}
        if backup.parent != bdir or backup.suffix != ".bak" or not backup.is_file():
            return {"ok": False, "error": "只能恢复自动备份里的文件"}
        meta = _read_backup_meta(backup)
        original = str(meta.get("original_path") or raw_path)
        if meta.get("original_path") and raw_path and not _same_resolved_path(original, raw_path):
            return {"ok": False, "error": "备份和文件不匹配"}
        p = Path(original).expanduser()
        if p.exists() and p.is_dir():
            return {"ok": False, "error": "目标是文件夹，不能用文件备份覆盖"}
        current_backup = _backup(p) if p.exists() else None
        p.parent.mkdir(parents=True, exist_ok=True)
        shutil.copy2(backup, p)
        return {
            "ok": True,
            "path": str(p),
            "backup_path": str(backup),
            "current_backup_path": current_backup,
        }
    except Exception as e:
        return {"ok": False, "error": str(e)[:120]}


# ────────────────────────────── 真 Agent 文件/命令工具的公共护栏 ──────────────────────────────

# 搜文件内容时只扫这些「文本类」扩展（没装 ripgrep 退回 Python os.walk 时用）；
# 二进制（图片/视频/Office 压缩包/可执行）跳过——搜内容没意义还容易乱码/卡死。
_TEXT_EXTS = {
    ".txt", ".md", ".markdown", ".csv", ".tsv", ".log", ".json", ".jsonl", ".yaml", ".yml",
    ".toml", ".ini", ".cfg", ".conf", ".env", ".py", ".js", ".ts", ".jsx", ".tsx", ".html",
    ".htm", ".css", ".scss", ".xml", ".sql", ".sh", ".bat", ".ps1", ".java", ".go", ".rs",
    ".c", ".cpp", ".h", ".hpp", ".php", ".rb", ".lua", ".vue", ".svelte", ".gitignore", "",
}
# os.walk 退回时单个文件扫描上限（字节）：超过当大文件跳过，防把几百 M 的日志读爆内存。
_GREP_MAX_FILE_BYTES = 5 * 1024 * 1024
# 危险命令黑名单（run_command 用）：命中任一即拒。覆盖删根/提权/裸盘写/格式化/递归放权/fork 炸弹/写系统要害区。
_DANGEROUS_PATTERNS = [
    r"\brm\b.*\s-[a-z]*[rf]",          # rm -rf / rm -fr / rm -r 等
    r"\bsudo\b",                        # 提权
    r"\bsu\b\s",                        # 切换用户
    r"\bdd\b.*\bif=",                   # dd if= 裸盘读写
    r"\bmkfs",                          # 格式化文件系统
    r"\bchmod\b.*-[a-z]*r[a-z]*\s*777",  # chmod -R 777 放权（命令已转小写，故匹配小写 r）
    r"\bchown\b.*-[a-z]*r",             # chown -R 递归改属主（同上）
    r":\(\)\s*\{.*\}\s*;\s*:",          # :(){ :|:& };: fork 炸弹
    r">\s*/dev/",                       # 写 /dev/*
    r"\b/dev/(sd|disk|null|zero|random)", # 触碰裸盘/设备
    r">\s*/boot",                       # 写 /boot
    r"/etc/ssh",                        # 动 ssh 配置/密钥
    r"\b(shutdown|reboot|halt|poweroff)\b",  # 关机重启
    r"\bmkfs|\bfdisk|\bparted",         # 分区/格式化
    r"\bcurl\b.*\|\s*(sh|bash)",        # curl ... | sh 远程执行
    r"\bwget\b.*\|\s*(sh|bash)",        # wget ... | sh
    # ── 数据外传命令（M5 ★2 堵注入→外传；注意 _check_command_safety 先 .lower() 再匹配） ──
    r"\bcurl\b.*\s-[a-z]*[dft]",        # curl -d/-F/-T 发送数据/上传文件（已转小写）
    r"\bcurl\b.*--data",                # curl --data/--data-binary/--data-urlencode
    r"\bcurl\b.*--upload-file",         # curl --upload-file
    r"\bcurl\b.*\s@",                   # curl @file 读本机文件发出
    r"\bscp\b",                         # scp 外传文件
    r"\bsftp\b",                        # sftp 外传文件
    r"\brsync\b",                       # rsync 同步文件到外部
    r"\bnc\b",                          # nc/netcat 网络外传
    r"\bncat\b",                        # ncat 网络外传
    r"\bnetcat\b",                      # netcat 网络外传
    r"\bwget\b.*--(post-data|post-file|body-data|body-file)",  # wget POST 外传
    r"\btftp\b",                        # tftp 文件传输
    r"\btelnet\b",                      # telnet 远程连接
]
# 禁用的 shell 操作符（防命令拼接/重定向/管道/子命令绕过黑名单）。命中即拒。
_SHELL_OPERATORS = ["&&", "||", "|", ";", ">", "<", "`", "$(", "&"]

# 环境变量泄漏专项：内置模型 key 打包在进程 env 里，裸读整包环境变量＝把 key 喂进模型上下文
# （模型看得到工具结果，读到 key 就等于泄露给了模型/prompt 注入）。只锚定「裸命令」形态——
# 整条命令(已 strip 过、无 shell 操作符残留)第一个词就是 env/printenv/set 本身：
#   - 会拦：`env`、`env -0`、`printenv`、`printenv PATH`、`set`、`set -x`
#   - 不误杀：`NODE_ENV=production node app.js`（第一个词是 NODE_ENV=..., 不是 env 本身）、
#     `node --env-file=.env server.js`（第一个词是 node）、`environment.sh`/`env_check.py`
#     这类以 env 开头但后面紧跟着单词字符的脚本名（\b 卡住，不会在词中间断开算命中）。
_ENV_LEAK_PATTERNS: list[tuple[str, str]] = [
    (r"^\s*env\b",
     "这条命令会读取整个环境变量表（env），本机进程的环境变量里存着内置模型的密钥——"
     "环境变量不能整包读出来喂给模型，出于安全拒绝执行。"),
    (r"^\s*printenv\b",
     "这条命令会读取环境变量（printenv），本机进程的环境变量里存着内置模型的密钥——"
     "环境变量不能整包读出来喂给模型，出于安全拒绝执行。"),
    (r"^\s*set\b",
     "这条命令会打印所有 shell 变量（set），本机进程的环境变量里存着内置模型的密钥——"
     "环境变量不能整包读出来喂给模型，出于安全拒绝执行。"),
]

# 上面的黑名单只能拦「裸命令」形态，`bash -c "env"`/`python3 -c "print(os.environ)"`/`/usr/bin/env`
# 套一层解释器就绕过了（复扫实测可达）。治本在下面：跑用户命令的子进程环境里直接剥掉密钥类变量——
# 用户的 shell 命令根本用不到模型 key，子进程里没有，怎么套壳都读不到。黑名单保留作第一道（报错话术友好）。
_SENSITIVE_ENV_KEY_PAT = re.compile(
    r"(API_KEY|APIKEY|SECRET|TOKEN|PASSWORD|PASSWD|PACK_KEY|PRIVATE_KEY|CREDENTIAL)", re.IGNORECASE
)


def sanitized_child_env() -> dict[str, str]:
    """给 run_command/run_background 子进程用的净化环境：剥掉密钥类变量，其余(PATH/HOME/LANG…)原样保留。"""
    return {k: v for k, v in os.environ.items() if not _SENSITIVE_ENV_KEY_PAT.search(k)}


def _resolve_dir(root_path: str, ctx=None) -> Path:
    """把一个【目录根】解析进沙箱并校验。语义同 _resolve（内容库/选定文件/全盘门控），
    但面向"在这个目录下找/搜/列"的场景。越界抛 ValueError。"""
    return _resolve(root_path, ctx)


# ────────────────────────────── 只读工具（无需审批） ──────────────────────────────

async def list_files(args: dict, ctx) -> str:
    """列文件。不传 path=列工作目录(没设则内容库)；传 path=列该目录（沙箱内随时、沙箱外要全盘）。"""
    raw = (args.get("path") or "").strip()
    if not raw:
        # 工作目录是当前会话的默认作业区。没设工作目录时，才回退到内容库。
        root = _resolve(".", ctx)
        if not root.exists():
            return f"当前工作目录不存在：{root}"
        if not root.is_dir():
            return f"当前工作目录不是文件夹：{root}"
        items = []
        for p in sorted(root.iterdir()):
            if ".backups" in p.parts:
                continue
            kind = "📁" if p.is_dir() else "📄"
            size = "" if p.is_dir() else f"  ({p.stat().st_size} 字节)"
            items.append(f"- {kind} {p.name}{size}")
        title = "当前工作目录" if (getattr(ctx, "working_dir", None) or "").strip() else "内容库文件"
        return f"{title} {root}：\n" + ("\n".join(items) if items else "（空目录）")
    # 传了 path：列指定目录（仅该层，不递归）
    try:
        d = _resolve_dir(raw, ctx)
    except ValueError as e:
        return f"列不了这个目录：{e}（沙箱外的目录需要老板开启「完全访问模式」）"
    if not d.exists():
        return f"目录不存在：{raw}"
    if not d.is_dir():
        return f"这不是目录（是个文件）：{raw}"
    items = []
    try:
        for p in sorted(d.iterdir()):
            kind = "📁" if p.is_dir() else "📄"
            size = "" if p.is_dir() else f"  ({p.stat().st_size} 字节)"
            items.append(f"- {kind} {p.name}{size}")
    except PermissionError:
        return f"没权限读这个目录：{raw}"
    return f"目录 {d}：\n" + ("\n".join(items) if items else "（空目录）")


async def find_files(args: dict, ctx) -> str:
    """按名字 glob 递归找文件（Claude Code 的 Glob 等价物）。
    args: root_path（绝对目录），pattern（glob 如 *.xlsx 或 **/采购*），max_results（默认 100）。只读。"""
    raw_root = (args.get("root_path") or "").strip()
    pattern = (args.get("pattern") or "").strip()
    if not raw_root:
        return "请给 root_path（要在哪个目录下找，绝对路径）。"
    if not pattern:
        return "请给 pattern（要找的文件名规则，glob，如 *.xlsx 或 **/采购*）。"
    try:
        max_results = int(args.get("max_results") or 100)
    except (TypeError, ValueError):
        max_results = 100
    max_results = max(1, min(max_results, 1000))
    try:
        root = _resolve_dir(raw_root, ctx)
    except ValueError as e:
        return f"找不了这个目录：{e}（沙箱外的目录需要老板开启「完全访问模式」）"
    if not root.exists():
        return f"目录不存在：{raw_root}"
    if not root.is_dir():
        return f"这不是目录：{raw_root}"
    # glob 默认只匹配当前层；想递归用户写 **/xxx。为贴合"递归按名字搜"的直觉，
    # 不含 / 的纯文件名模式（如 *.xlsx）自动当成 **/ 递归找。
    glob_pat = pattern if "/" in pattern else f"**/{pattern}"
    hits: list[Path] = []
    total = 0
    try:
        for p in root.glob(glob_pat):
            if p.is_file() and ".backups" not in p.parts:
                total += 1
                if len(hits) < max_results:
                    hits.append(p)
    except (ValueError, OSError) as e:
        return f"搜的时候出错了：{e}"
    if total == 0:
        return f"在 {root} 下没找到匹配「{pattern}」的文件。"
    lines = [f"在 {root} 下找到 {total} 个匹配「{pattern}」的文件" +
             (f"（只列前 {max_results} 个）：" if total > len(hits) else "：")]
    lines += [f"- {p}" for p in hits]
    return "\n".join(lines)


def _grep_with_ripgrep(root: Path, query: str, max_results: int) -> list[str] | None:
    """有 rg 就用它搜内容（快、自动跳二进制/.gitignore）。返回 "文件:行号:命中行" 列表；
    没装 rg 返回 None（让调用方退回 Python 实现）。"""
    rg = shutil.which("rg")
    if not rg:
        return None
    try:
        # -n 行号 -H 带文件名 -e 把 query 当独立模式（防 query 以 - 开头被当选项）-m 每文件命中上限
        proc = subprocess.run(
            [rg, "-n", "-H", "--max-count", "50", "-e", query, str(root)],
            shell=False, capture_output=True, text=True, timeout=30,
        )
    except (subprocess.TimeoutExpired, OSError):
        return None
    out_lines = (proc.stdout or "").splitlines()
    results: list[str] = []
    for ln in out_lines:
        results.append(ln)
        if len(results) >= max_results:
            break
    return results


def _grep_with_python(root: Path, query: str, max_results: int) -> list[str]:
    """没装 rg 的退路：os.walk + 正则，只扫文本类扩展、跳大文件/二进制。"""
    try:
        rx = re.compile(query)
    except re.error:
        # query 不是合法正则 → 当普通子串搜（转义）
        rx = re.compile(re.escape(query))
    results: list[str] = []
    for dirpath, dirnames, filenames in os.walk(root):
        # 跳过备份目录与隐藏的版本控制目录，省时间
        dirnames[:] = [d for d in dirnames if d not in (".backups", ".git", "node_modules", "__pycache__")]
        for fn in filenames:
            fp = Path(dirpath) / fn
            if fp.suffix.lower() not in _TEXT_EXTS:
                continue
            try:
                if fp.stat().st_size > _GREP_MAX_FILE_BYTES:
                    continue
            except OSError:
                continue
            try:
                with fp.open("r", encoding="utf-8", errors="ignore") as f:
                    for i, line in enumerate(f, 1):
                        if rx.search(line):
                            results.append(f"{fp}:{i}:{line.rstrip()[:300]}")
                            if len(results) >= max_results:
                                return results
            except (OSError, UnicodeError):
                continue
    return results


async def search_in_files(args: dict, ctx) -> str:
    """在目录下按内容搜文件（Claude Code 的 Grep 等价物）。优先 ripgrep，没装退回 Python。
    args: root_path（绝对目录），query（关键词/正则），max_results（默认 100）。只读。"""
    raw_root = (args.get("root_path") or "").strip()
    query = args.get("query") or ""
    if not raw_root:
        return "请给 root_path（要在哪个目录下搜，绝对路径）。"
    if not query.strip():
        return "请给 query（要搜的关键词或正则）。"
    try:
        max_results = int(args.get("max_results") or 100)
    except (TypeError, ValueError):
        max_results = 100
    max_results = max(1, min(max_results, 1000))
    try:
        root = _resolve_dir(raw_root, ctx)
    except ValueError as e:
        return f"搜不了这个目录：{e}（沙箱外的目录需要老板开启「完全访问模式」）"
    if not root.exists():
        return f"目录不存在：{raw_root}"
    if not root.is_dir():
        return f"这不是目录：{raw_root}"
    results = _grep_with_ripgrep(root, query, max_results)
    if results is None:
        results = _grep_with_python(root, query, max_results)
    if not results:
        return f"在 {root} 下没有内容含「{query}」的文件。"
    head = f"在 {root} 下搜到「{query}」（文件:行号:命中行）："
    return head + "\n" + "\n".join(results)


# ────────────────────────────── 跑命令（最险·硬门控全盘 + 黑名单 + 禁操作符 + 审批闸） ──────────────────────────────

def _check_command_safety(command: str) -> str | None:
    """命令安全检查。安全返回 None；不安全返回【拒绝原因】（人话）。
    顺序：① 必须开全盘（由 handler 在调用前先判，这里不重复）② 禁 shell 操作符 ③ 环境变量泄漏专项 ④ 危险黑名单。"""
    cmd = (command or "").strip()
    if not cmd:
        return "命令是空的，没东西可跑。"
    # ② 禁 shell 操作符（防拼接/重定向/管道/子命令绕过黑名单）
    for op in _SHELL_OPERATORS:
        if op in cmd:
            return (f"命令里有 shell 操作符「{op}」，不允许（防止命令拼接/重定向绕过安全检查）。"
                    "请一次只跑一条简单命令，不要用 && || | ; > < 反引号 $() 这些。")
    low = cmd.lower()
    # ③ 环境变量泄漏专项（裸 env/printenv/set：内置模型 key 在进程 env 里，读出来就是喂给模型）
    for pat, reason in _ENV_LEAK_PATTERNS:
        if re.search(pat, low):
            return reason
    # ④ 危险黑名单
    for pat in _DANGEROUS_PATTERNS:
        if re.search(pat, low):
            return f"这条命令命中了危险操作黑名单（{pat}），出于安全拒绝执行。"
    return None


def _run_command_fatal_reason(args: dict, ctx) -> str | None:
    """审批闸 2.0 · ①"撞墙才问"的命令类判据（挂到 run_command 的 `Tool.fatal_reason_for`）。

    判据跟 handler 内部 `_check_command_safety` 完全一致（同一个函数，不新造一套黑名单）——
    命中的都是【致命/永不批】类（删根/提权/裸盘/格式化/fork炸弹/关机/远程执行/外传数据等，
    见 `_DANGEROUS_PATTERNS`/`_ENV_LEAK_PATTERNS`/`_SHELL_OPERATORS`），黑名单优先级最高、
    任何权限档都不给开口子（这三份清单里的模式经安全评审后判定"没有正当理由需要 Agent 自主跑"，
    故意不设"可挽救、批准后放行"的中间档——这不是本单新引入的收紧，只是把原本【approval 卡通过后
    handler 才二次拒绝】的既有行为，挪到 `_plan_tool_call` 更早的地方直接说清楚原因，
    别让老板在 ask/auto_files 档看一张"确认执行"卡、点了也没用，也别让 full 档看着"已自动执行"、
    结果其实是执行时才失败的糊涂账。"""
    command = (args or {}).get("command") or ""
    return _check_command_safety(command)


# 审批闸 2.0 · ③ Roo 式安全前缀白名单：只读、无副作用的命令前缀，最长前缀匹配命中 → 任何权限档
# （含"逐项确认"）都可以自动放行、不用弹卡。刻意保持这份清单短小、每条都明显无副作用——
# 白名单越大风险面越大，宁可漏放（仍会走既有"跑权限档"的普通审批流程）也不错放。
_SAFE_COMMAND_PREFIXES: tuple[str, ...] = (
    "ls", "cat", "pwd", "echo", "head", "tail",
    "git status", "git log", "git diff", "git branch",
    "whoami", "date", "wc", "file", "stat", "df", "du", "which", "uname",
)


def _matches_safe_prefix(command: str) -> bool:
    """最长前缀匹配：命令原文（原样或后跟一个空格）等于清单里的某条 → 命中。
    按长度从长到短比对，防"git"这种裸前缀（本清单没有，注释留痕）误吞"git status"之外的危险子命令。"""
    cmd = (command or "").strip().lower()
    if not cmd:
        return False
    for prefix in sorted(_SAFE_COMMAND_PREFIXES, key=len, reverse=True):
        if cmd == prefix or cmd.startswith(prefix + " "):
            return True
    return False


# cat/head/tail 会把文件【内容】读进模型上下文——即便命令前缀在白名单里，也不能豁免过一遍
# read_file 同款的敏感文件识别（_is_sensitive_file），否则 "cat ~/.ssh/id_rsa" 这种就绕开了
# read_file 本该有的确认闸、零点击把密钥灌进模型上下文。ls/pwd/echo/git status 等不读文件内容，
# 不受此限（只读文件名/状态列表，不算"内容"）。
_CONTENT_READING_PREFIXES = ("cat", "head", "tail")


def _run_command_safe_prefix(args: dict, ctx) -> bool:
    """挂到 run_command 的 `Tool.safe_prefix_for`：命令命中安全前缀白名单 → 免弹审批卡直接跑。
    先复核一遍 `_check_command_safety`（双保险，理论上安全前缀不会同时命中危险黑名单，
    但判定壳子出错/清单以后被改坏时，危险黑名单必须优先赢）；
    再对 cat/head/tail 这类"读文件内容"的命令，逐个参数过一遍 `_is_sensitive_file`——
    命中就不豁免（退回常规审批流程，不是拒绝），防止零点击读敏感文件。"""
    command = (args or {}).get("command") or ""
    if _check_command_safety(command) is not None:
        return False
    if not _matches_safe_prefix(command):
        return False
    cmd = command.strip().lower()
    if any(cmd == p or cmd.startswith(p + " ") for p in _CONTENT_READING_PREFIXES):
        try:
            parts = shlex.split(command)
        except ValueError:
            return False  # 参数解析不了（引号没配对等）→ 保守不豁免，走常规审批
        for token in parts[1:]:
            if not token.startswith("-") and _is_sensitive_file(token):
                return False
    return True


def _kill_proc_group(proc, posix: bool) -> None:
    """超时掐断：POSIX 下整组 SIGKILL（连子进程 fork 出的孙子进程一起杀，不留孤儿占端口/吃 CPU）；
    取不到组（已退出/无权限）或非 POSIX → 退回只杀直接子进程。绝不抛。"""
    if proc is None or getattr(proc, "pid", None) is None:
        return
    try:
        if posix:
            os.killpg(os.getpgid(proc.pid), signal.SIGKILL)
            return
    except (ProcessLookupError, PermissionError, OSError):
        pass
    try:
        proc.kill()
    except ProcessLookupError:
        pass


async def run_command(args: dict, ctx) -> str:
    """在老板本机跑一条命令（Claude Code 的 Bash 等价物·最危险）。
    硬门控：必须开「完全访问模式」(full_disk_access)；禁 shell 操作符；危险黑名单；shell=False；超时；输出截断。
    args: command（字符串），timeout_sec（默认 30），cwd（工作目录，可选）。审批跟权限档走（L1/L2弹卡、L3自己跑）；危险黑名单与档位无关、永远拦死。"""
    # ① 硬门控：没开完全访问模式 → 直接拒（命令工具风险最高，绝不在普通模式下放行）
    if not getattr(ctx, "full_disk_access", False):
        return "跑命令需要先开启「完全访问模式」。普通模式下我只能动内容库和你选定的文件，不能跑命令。"
    command = (args.get("command") or "").strip()
    reason = _check_command_safety(command)
    if reason:
        return f"拒绝执行：{reason}"
    try:
        timeout_sec = float(args.get("timeout_sec") or 30)
    except (TypeError, ValueError):
        timeout_sec = 30.0
    timeout_sec = max(1.0, min(timeout_sec, 300.0))
    cwd = (args.get("cwd") or "").strip() or (getattr(ctx, "working_dir", None) or None)
    if cwd and not Path(cwd).is_dir():
        return f"工作目录不存在：{cwd}"
    try:
        parts = shlex.split(command)
    except ValueError as e:
        return f"命令解析失败（引号没配对？）：{e}"
    if not parts:
        return "命令是空的，没东西可跑。"
    # 命令边跑边显示（对标 Claude Code）：用 asyncio 子进程逐行读 stdout/stderr，
    # 若 ctx.progress_emit 在则每段实时推出（流式循环据此 yield tool_progress）；同时攒全量供最终结果。
    emit = getattr(ctx, "progress_emit", None)
    stdout_parts: list[str] = []
    stderr_parts: list[str] = []

    async def _pump(stream, sink, kind):
        if stream is None:
            return
        while True:
            raw = await stream.readline()
            if not raw:
                break
            line = raw.decode("utf-8", errors="replace")
            sink.append(line)
            if emit:
                try:
                    emit({"type": "tool_progress", "tool": "run_command", "stream": kind, "chunk": line})
                except Exception:
                    pass

    _posix = (os.name == "posix")
    try:
        proc = await asyncio.create_subprocess_exec(
            *parts, stdout=asyncio.subprocess.PIPE, stderr=asyncio.subprocess.PIPE, cwd=cwd,
            # POSIX：子进程自成会话/进程组(setsid) → 超时能整组掐断，连它 fork 的孙子进程一起杀，不留孤儿。
            # Windows 忽略此参数（无害），仍走下方 proc.kill() 兜底。
            start_new_session=_posix,
            # 密钥类环境变量不给子进程：黑名单只能拦裸 `env`，`bash -c "env"`/`python3 -c "os.environ"`
            # 一包一层就绕过了。治本 = 子进程环境里压根没有 key（用户命令也用不到模型 key）。
            env=sanitized_child_env(),
        )
    except FileNotFoundError:
        return f"找不到这个命令（没装或不在 PATH 里）：{parts[0]}"
    except OSError as e:
        return f"命令没跑起来：{e}"

    try:
        await asyncio.wait_for(
            asyncio.gather(
                _pump(proc.stdout, stdout_parts, "stdout"),
                _pump(proc.stderr, stderr_parts, "stderr"),
                proc.wait(),
            ),
            timeout=timeout_sec,
        )
    except asyncio.TimeoutError:
        _kill_proc_group(proc, _posix)
        try:  # 回收已杀进程，别留僵尸
            await asyncio.wait_for(proc.wait(), timeout=3)
        except (asyncio.TimeoutError, ProcessLookupError, Exception):
            pass
        return f"命令跑超时了（超过 {int(timeout_sec)} 秒自动掐断）：{command}"

    stdout_show = _truncate_output("".join(stdout_parts))
    stderr_show = _truncate_output("".join(stderr_parts))
    out = [f"命令：{command}", f"返回码：{proc.returncode}"]
    if stdout_show:
        out.append(f"【标准输出】\n{stdout_show}")
    if stderr_show:
        out.append(f"【错误输出】\n{stderr_show}")
    if not stdout_show and not stderr_show:
        out.append("（无输出）")
    return "\n".join(out)


_TRUNC_MAX_LINE_CHARS = 4000  # 单行字符护栏：防一行塞几十万字符（压成一行的 JSON/日志）把输出撑爆


def _truncate_output(text: str, head_lines: int = 30, tail_lines: int = 70) -> str:
    """命令输出按行截断：保头 head_lines 行 + 保尾 tail_lines 行，中间插一行省略提示。

    2-2 修复：旧版"只留前 100 行"会把报错切没——脚本/测试失败信息几乎总在输出末尾，模型看到
    开头一片正常就误判"跑成了"。改成头尾各留一截（Claude Code 的截中间保头尾做法），报错保留得住。
    另加单行字符护栏：任何一行超 _TRUNC_MAX_LINE_CHARS 字符也再截一道，防单行内容本身就撑爆上下文。
    """
    if not text:
        return ""
    lines = text.splitlines()
    total = len(lines)
    max_lines = head_lines + tail_lines
    if total <= max_lines:
        kept = lines
        omitted = 0
    else:
        omitted = total - max_lines
        kept = lines[:head_lines] + [f"…（中间省略 {omitted} 行）…"] + lines[-tail_lines:]
    capped = [
        (ln if len(ln) <= _TRUNC_MAX_LINE_CHARS
         else f"{ln[:_TRUNC_MAX_LINE_CHARS]}…[本行过长共 {len(ln)} 字符，已截断]")
        for ln in kept
    ]
    result = "\n".join(capped).rstrip()
    if omitted:
        result += f"\n…[输出共 {total} 行，已保留头 {head_lines} 行 + 尾 {tail_lines} 行]"
    return result


_DOC_MAX_CHARS = 9000  # PDF/Word/PPT 提取文字封顶，防撑爆上下文；超了给清晰提示


def _cap_doc(text: str, kind: str) -> str:
    if len(text) > _DOC_MAX_CHARS:
        return text[:_DOC_MAX_CHARS] + f"\n\n…[{kind}较长，已读前 {_DOC_MAX_CHARS} 字；要看后面或某部分请说明]"
    return text


def _read_pdf(path) -> str:
    try:
        from pypdf import PdfReader
        reader = PdfReader(str(path))
        parts = []
        for i, page in enumerate(reader.pages, 1):
            t = (page.extract_text() or "").strip()
            if t:
                parts.append(f"--- 第 {i} 页 ---\n{t}")
        if not parts:
            return "（这份 PDF 没提取到文字——多半是扫描件/图片型 PDF，要 OCR 才能读）"
        return _cap_doc("\n\n".join(parts), "PDF")
    except Exception as e:  # noqa: BLE001
        return f"（PDF 读取失败：{e}）"


def _read_docx(path) -> str:
    try:
        import docx
        doc = docx.Document(str(path))
        parts = [p.text for p in doc.paragraphs if p.text.strip()]
        for tbl in doc.tables:
            for row in tbl.rows:
                cells = [c.text.strip() for c in row.cells]
                if any(cells):
                    parts.append(" | ".join(cells))
        if not parts:
            return "（这份 Word 文档没有文字内容）"
        return _cap_doc("\n".join(parts), "Word 文档")
    except Exception as e:  # noqa: BLE001
        return f"（Word 读取失败：{e}）"


def _read_pptx(path) -> str:
    try:
        from pptx import Presentation
        prs = Presentation(str(path))
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = [s.text_frame.text.strip() for s in slide.shapes
                     if s.has_text_frame and s.text_frame.text.strip()]
            if texts:
                parts.append(f"--- 第 {i} 页 ---\n" + "\n".join(texts))
        if not parts:
            return "（这份 PPT 没有文字内容）"
        return _cap_doc("\n\n".join(parts), "PPT")
    except Exception as e:  # noqa: BLE001
        return f"（PPT 读取失败：{e}）"


async def read_file(args: dict, ctx) -> str:
    """读一个文件的内容，给 Agent 看（编辑前先读）。文本直接读；Excel 列非空单元格；PDF/Word/PPT 提取文字。"""
    path = _resolve(args["path"], ctx)
    if not path.exists():
        return f"文件不存在：{args['path']}"
    if path.suffix.lower() in (".xlsx", ".xlsm"):
        from openpyxl import load_workbook
        wb = load_workbook(path, data_only=True)
        # read_file 是 read_only：超长结果走硬截断、不落盘。大报表若逐格全列，会被通用截断把单元格切碎、
        # 语义全毁。故在此先按非空单元格数封顶，只给前若干行的完整内容 + 表尺寸提示，保证可读。
        MAX_CELLS = 1200
        lines = []
        emitted = 0
        truncated = False
        for ws in wb.worksheets:
            lines.append(f"# 工作表「{ws.title}」（{ws.max_row} 行 × {ws.max_column} 列）")
            if truncated:
                lines.append("…[已达读取上限，本表略]")
                continue
            for row in ws.iter_rows():
                if emitted >= MAX_CELLS:
                    truncated = True
                    lines.append(f"…[表太大，已读约 {MAX_CELLS} 个非空单元格、余下略；要看特定区域请说明行列范围]")
                    break
                for cell in row:
                    if cell.value is not None:
                        lines.append(f"{cell.coordinate}={cell.value!r}")
                        emitted += 1
        return "\n".join(lines) if lines else "（空表）"
    if path.suffix.lower() == ".pdf":
        return _read_pdf(path)
    if path.suffix.lower() == ".docx":
        return _read_docx(path)
    if path.suffix.lower() == ".pptx":
        return _read_pptx(path)
    # 缺口 F：图片不是"二进制不便读取"——把路径挂进回灌队列(ctx.pending_view_images)，
    # loop 在本批 tool 结果配对完整后经 _drain_view_images 拼成一条 user 图片消息注入(走带原始尺寸
    # 标签的 multimodal 通道、自动缩到 <=1568)，模型【下一轮】就真看见这张图。对标官方 FileReadTool
    # 读图自动转 image block 回灌。沙箱/敏感闸已在到这之前走过(_resolve + read_file 的 requires_approval_for)。
    from services.agent.multimodal import is_image
    if is_image(path):
        pending = getattr(ctx, "pending_view_images", None)
        if pending is not None:
            try:
                pending.append(str(path))
            except Exception:  # noqa: BLE001 — 回灌是尽力而为，绝不让读图崩掉工具
                pass
        return f"已读取图片〈{path.name}〉，将在下一轮直接看到它的画面（据此判断与继续）。"
    try:
        return path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        return f"（二进制文件，{path.stat().st_size} 字节，不便直接读取）"


# ────────────────────────────── 写/改工具（走审批闸） ──────────────────────────────

async def write_file(args: dict, ctx) -> str:
    """把内容写到内容库里的一个文件（新建或覆盖）。args: path, content。覆盖前自动备份。"""
    path = _resolve(args["path"], ctx)
    path.parent.mkdir(parents=True, exist_ok=True)
    backup = _backup(path)
    path.write_text(args["content"], encoding="utf-8")
    msg = f"已写入 {path.name}（{len(args['content'])} 字）。"
    if backup:
        msg += f" 原件已备份。"
    return msg


_EDIT_RECEIPT_CONTEXT_LINES = 3   # 改动位置前后各留几行，让模型能看清周边代码没被误伤
_EDIT_RECEIPT_MAX_LINES = 200     # 片段太大（整段替换几百行）再兜底截，防回执本身撑爆上下文


def _numbered_snippet(full_text: str, start_line: int, end_line: int, context: int) -> str:
    """按 1-based 行号取 [start_line, end_line] 及前后各 context 行，拼成带行号的片段（新内容视角）。"""
    lines = full_text.splitlines()
    total = len(lines)
    lo = max(1, start_line - context)
    hi = min(total, end_line + context)
    if hi - lo + 1 > _EDIT_RECEIPT_MAX_LINES:
        half = _EDIT_RECEIPT_MAX_LINES // 2
        head = [f"{i:>5}\t{lines[i - 1]}" for i in range(lo, lo + half)]
        tail = [f"{i:>5}\t{lines[i - 1]}" for i in range(hi - half + 1, hi + 1)]
        return "\n".join(head) + f"\n      …（片段太长，中间省略 {hi - lo + 1 - _EDIT_RECEIPT_MAX_LINES} 行）…\n" + "\n".join(tail)
    return "\n".join(f"{i:>5}\t{lines[i - 1]}" for i in range(lo, hi + 1))


async def edit_file(args: dict, ctx) -> str:
    """改文本文件的某一段：把 old_text 精确替换成 new_text（同我改代码的方式）。改前备份。
    args: path, old_text, new_text。"""
    path = _resolve(args["path"], ctx)
    if not path.exists():
        return f"文件不存在：{args['path']}"
    text = path.read_text(encoding="utf-8")
    old, new = args["old_text"], args["new_text"]
    n = text.count(old)
    if n == 0:
        return f"没找到要替换的内容，未改动。"
    if n > 1:
        return f"要替换的内容出现 {n} 次（不唯一），为安全未改动；请给更具体的上下文。"
    idx = text.find(old)
    start_line = text.count("\n", 0, idx) + 1
    new_full = text.replace(old, new)
    end_line = start_line + new.count("\n")
    backup = _backup(path)
    path.write_text(new_full, encoding="utf-8")
    snippet = _numbered_snippet(new_full, start_line, end_line, _EDIT_RECEIPT_CONTEXT_LINES)
    where = f"第 {start_line} 行" if end_line == start_line else f"第 {start_line}-{end_line} 行"
    receipt = f"已修改 {path.name}（{where}附近，改动后的样子）：\n{snippet}"
    if backup:
        receipt += "\n（原件已备份，可回滚）"
    return receipt


async def edit_excel(args: dict, ctx) -> str:
    """直接改 Excel 报表的单元格（改营业额、加一列提成等）。改前备份、改后回传逐格 diff。
    args: path, changes=[{cell:'B2', value:8600}, ...]（cell 用 A1 式坐标；多表加 sheet）。"""
    from openpyxl import load_workbook
    path = _resolve(args["path"], ctx)
    if not path.exists():
        return f"报表不存在：{args['path']}"
    wb = load_workbook(path)
    diffs = []
    for ch in args.get("changes", []):
        ws = wb[ch["sheet"]] if ch.get("sheet") else wb.active
        cell = ch["cell"]
        try:
            old = ws[cell].value
            ws[cell] = ch["value"]
        except AttributeError:  # 合并单元格的非左上角格只读
            return f"{cell} 是合并单元格，改不了——请改合并区左上角那个格子。"
        except (ValueError, KeyError):
            return f"坐标无效：{cell}。请用 A1 式坐标（如 B2）。"
        diffs.append(f"{ws.title}!{cell}: {old!r} → {ch['value']!r}")
    backup = _backup(path)  # 改成功才备份，避免改失败留孤儿备份
    wb.save(path)
    return f"已改 {path.name}：\n" + "\n".join(diffs) + ("\n（原件已备份，可回滚）" if backup else "")


async def delete_file(args: dict, ctx) -> str:
    """删除内容库里的一个文件（删前自动备份到 .backups，可恢复）。args: path。
    高危不可逆动作——force_confirm=True，任何权限档都强制弹确认后才会执行（不被"跳过确认"旁路）。
    只删单个文件、不删文件夹（删目录风险大，挡住）。"""
    path = _resolve(args["path"], ctx)
    if not path.exists():
        return f"文件不存在：{args['path']}（无需删除）"
    if path.is_dir():
        return f"「{path.name}」是文件夹——为安全本工具只删单个文件，不删目录。请逐个文件删。"
    backup = _backup(path)
    path.unlink()
    msg = f"已删除 {path.name}。"
    if backup:
        msg += " 原件已备份到 .backups，可恢复。"
    return msg


# ────────────────────────────── 召回：翻老板本机攒下的历史内容（真 RAG·语义检索） ──────────────────────────────

async def recall_my_content(args, ctx) -> str:
    """语义检索老板过去生成的内容（"找我上次那条…""跟之前类似的"）。先惰性补建索引再搜。"""
    store = getattr(ctx, "store", None)
    if store is None or getattr(store, "id", None) is None:
        return "（拿不到当前门店，没法翻历史。）"
    from services.rag.recall import backfill_from_generations, recall
    await backfill_from_generations(ctx.db, store.id)
    hits = recall(str(store.id), args.get("query", "") or "", top=5)
    if not hits:
        return "没在你过去的内容里找到相关的。要不直接说需求，我现写一条。"
    lines = []
    for h in hits:
        snippet = " ".join((h["text"] or "").split())[:200]
        lines.append(f"- {snippet}")
    return "翻到这些你以前写过的相关内容（可参考/在此基础上改）：\n" + "\n".join(lines)


# ────────────────────────────── POS 真诊断：读老板导出的报表 → 基于真实数字诊断 ──────────────────────────────

async def diagnose_from_pos(args: dict, ctx) -> str:
    """读老板从收银系统导出的报表(Excel)，喂进经营诊断引擎(决策树+指标库)，给【有数有据】的诊断——
    不再凭老板一句口述泛泛而谈。桌面独有：报表在老板自己电脑上，要他先选定文件。"""
    file = (args.get("file") or "").strip()
    # 没给路径 / 给的找不到 → 兜底用老板【当场选定的报表】(.xlsx)，省得大脑非得报对完整路径
    def _picked_report() -> str | None:
        for p in (getattr(ctx, "allowed_paths", None) or []):
            ps = str(p)
            if ps.lower().endswith((".xlsx", ".xls")) and Path(ps).exists():
                return ps
        return None
    if not file:
        file = _picked_report() or ""
    if not file:
        return "请先用文件选择器选一下你从收银系统导出的报表（.xlsx），我照着真实数据帮你看。"
    path = _resolve(file, ctx)
    if not path.exists():
        fallback = _picked_report()
        if fallback:
            path = Path(fallback)
        else:
            return f"没找到这个文件：{file}。麻烦用文件选择器重新选一下导出的报表。"
    from services.report_reader import extract_report_indicators
    _indicators, summary = extract_report_indicators(str(path))
    if _indicators:
        try:
            from services.agent.tools import _memory_candidate_from_report_summary
            from services.memory_service import add_pending_memory_candidate
            candidate = _memory_candidate_from_report_summary(summary)
            if candidate and getattr(ctx, "db", None) is not None and getattr(ctx, "store", None) is not None:
                await add_pending_memory_candidate(ctx.db, ctx.store.id, candidate, "operational")
        except Exception:
            logger.debug("failed to add pending memory from POS report", exc_info=True)
    # 复用 read_file 的读法（Excel 列出非空单元格）拿到真实数字。
    # ⚠️ 必须读【解析/兜底后的 path】，不是原始 file——否则给的路径不存在时会把"文件不存在"当数据喂给诊断引擎。
    data_text = summary if _indicators else await read_file({"path": str(path)}, ctx)
    situation = (
        "以下是这家店从收银系统导出的【真实经营数据】。请**基于这些具体数字**诊断："
        "引用关键数字、算出关键比率(如台费占比/空台时段)、指出异常项，再给可落地建议，别泛泛而谈：\n"
        f"{data_text}"
    )
    focus = (args.get("focus") or "").strip()
    if focus:
        situation += f"\n\n老板想重点看：{focus}"
    from services.diagnosis_service import analyze_diagnosis  # 懒导入，避免顶层耦合
    gen = await analyze_diagnosis(
        ctx.db, ctx.store, getattr(ctx, "user", None),
        problem_area=(args.get("problem_area") or "revenue"),
        current_situation=situation,
    )
    return gen.result


# ────────────────────────────── M5b 敏感文件读取确认闸 ──────────────────────────────

_SENSITIVE_EXACT_NAMES: set[str] = {
    ".env", ".netrc", ".git-credentials", ".npmrc", ".pgpass", "credentials",
    "id_rsa", "id_dsa", "id_ecdsa", "id_ed25519",
}

_SENSITIVE_EXTS: set[str] = {
    ".pem", ".key", ".p12", ".pfx", ".keystore", ".jks", ".kdbx", ".ovpn",
}

_SENSITIVE_DIR_SEGMENTS: set[str] = {
    ".ssh", ".aws", ".gnupg", ".kube",
}

_BROWSER_LOGIN_DBS: set[str] = {
    "login data", "key4.db", "logins.json", "cookies",
}


def _is_sensitive_file(path: str | None) -> bool:
    """判断路径是否指向可能含密钥/凭据的敏感文件（大小写不敏感）。"""
    if not path or not path.strip():
        return False
    normalized = path.strip().replace("\\", "/").lower()
    parts = normalized.split("/")
    basename = parts[-1]
    if not basename:
        return False
    if basename in _SENSITIVE_EXACT_NAMES:
        return True
    if basename.startswith(".env."):
        return True
    dot = basename.rfind(".")
    if dot > 0 and basename[dot:] in _SENSITIVE_EXTS:
        return True
    dirs = parts[:-1]
    for seg in dirs:
        if seg in _SENSITIVE_DIR_SEGMENTS:
            return True
    joined = "/" + "/".join(parts) + "/"
    if "/.config/gcloud/" in joined:
        return True
    if basename == "config.json" and len(parts) >= 2 and parts[-2] == ".docker":
        return True
    if "/library/keychains/" in joined:
        return True
    if basename in _BROWSER_LOGIN_DBS:
        return True
    return False


def _sensitive_read_reason(args: dict, ctx) -> dict:
    name = Path(args.get("path", "?") or "?").name
    return {
        "what": f"读取文件「{name}」",
        "why": "该文件可能含密钥/密码/凭据等敏感信息。读取后内容会发给 AI 模型处理，存在泄露风险。",
        "impact": "确认后才会读取文件内容；拒绝则跳过、不读。文件本身不会被修改。",
    }


def _sensitive_read_preview(args: dict, ctx) -> str:
    path = args.get("path", "?")
    return f"将读取可能含敏感信息的文件「{path}」。读取后内容会发给 AI 处理。"


# ────────────────────────────── 审批预览（确认前给老板看"会改成什么"，不再瞎确认） ──────────────────────────────

def _name_of(args: dict) -> str:
    return Path(args.get("path", "?") or "?").name


def preview_edit_excel(args: dict, ctx) -> str:
    """改 Excel 前的人话 diff：逐格 旧值→新值（读现值算）。读不到就只列要写的新值，绝不抛错。"""
    changes = args.get("changes", []) or []
    lines = [f"改报表《{_name_of(args)}》，共 {len(changes)} 处："]
    wb = None
    try:
        from openpyxl import load_workbook
        path = _resolve(args["path"], ctx)
        if path.exists():
            wb = load_workbook(path, data_only=True)
    except Exception:
        wb = None
    for ch in changes:
        cell, sheet = ch.get("cell", "?"), ch.get("sheet")
        old = "?"
        try:
            if wb is not None:
                ws = wb[sheet] if sheet else wb.active
                old = ws[cell].value
        except Exception:
            old = "?"
        loc = f"{sheet}!{cell}" if sheet else cell
        lines.append(f"  {loc}：{old!r} → {ch.get('value')!r}")
    return "\n".join(lines)


def preview_edit_file(args: dict, ctx) -> str:
    old = (args.get("old_text") or "")[:140]
    new = (args.get("new_text") or "")[:140]
    return f"改文件《{_name_of(args)}》：\n- 原：{old}\n+ 改：{new}"


def preview_run_command(args: dict, ctx) -> str:
    """跑命令前给老板看的预览：把【将执行的命令原文】清清楚楚展示，让他看原文再点确认。"""
    command = (args.get("command") or "").strip() or "（空）"
    cwd = (args.get("cwd") or "").strip()
    try:
        timeout_sec = int(float(args.get("timeout_sec") or 30))
    except (TypeError, ValueError):
        timeout_sec = 30
    lines = ["将在你电脑上执行命令（看清原文再确认）：", f"  $ {command}"]
    if cwd:
        lines.append(f"  工作目录：{cwd}")
    lines.append(f"  超时：{timeout_sec} 秒")
    return "\n".join(lines)


def preview_write_file(args: dict, ctx) -> str:
    content = args.get("content") or ""
    exists = False
    try:
        exists = _resolve(args["path"], ctx).exists()
    except Exception:
        exists = False
    snippet = content[:200] + ("…" if len(content) > 200 else "")
    return f"{'覆盖' if exists else '新建'}文件《{_name_of(args)}》（{len(content)} 字）：\n{snippet}"


def preview_delete_file(args: dict, ctx) -> str:
    """删文件前给老板看的预览：清楚说明要删哪个文件、删了能不能恢复。"""
    name = _name_of(args)
    info = ""
    try:
        p = _resolve(args["path"], ctx)
        if not p.exists():
            info = "（文件不存在，无需删）"
        elif p.is_dir():
            info = "（这是文件夹，本工具不删目录）"
        else:
            info = f"（{p.stat().st_size} 字节）"
    except Exception:
        info = ""
    return f"⚠️ 将【删除】文件《{name}》{info}。删前自动备份到 .backups、可恢复；你确认后才真正删。"


# ────────────────────────────── 工具定义（人看得懂的描述，大脑据此选） ──────────────────────────────

_LOCAL_TOOLS = [
    Tool(
        name="recall_my_content",
        description="检索老板【以前生成过】的内容（按意思找，不是按关键词）。当老板说"
                    "『找我上次那条…』『跟之前类似的』『把以前效果好的那条改改』『我之前写过的XX』时调用——"
                    "先翻历史找出相关的几条，再据此改写/参考，比从零写更贴老板的风格。",
        parameters={"type": "object", "properties": {
            "query": {"type": "string", "description": "要找的内容/主题，原话即可，如'双十一活动朋友圈'"},
        }, "required": ["query"]},
        handler=recall_my_content,
        read_only=True,
    ),
    Tool(
        name="diagnose_from_pos",
        description="读老板从【收银系统导出的报表 Excel】(营业额/台时/各项收入/上钟等)，基于真实数字做经营诊断。"
                    "当老板说『我导出了数据你帮我看看 / 看看这个月经营 / 分析下这张报表 / 照着我的真实数据诊断』"
                    "并且选定了一个报表文件时调用——比凭口述诊断准得多、会引用具体数字。",
        parameters={"type": "object", "properties": {
            "file": {"type": "string", "description": "老板选定的 POS 导出报表文件路径(.xlsx)"},
            "focus": {"type": "string", "description": "想重点看什么(可选)，如'为什么周二下午营收低'"},
            "problem_area": {"type": "string", "description": "问题领域(可选)：revenue/traffic/customer_loss/staff/competition/off_season"},
        }, "required": ["file"]},
        handler=diagnose_from_pos,
        read_only=True,
    ),
    Tool(
        name="list_files",
        description="列文件。不传 path＝列当前工作目录（没选工作目录时才列本机内容库）；"
                    "传 path＝列那个目录里有什么（沙箱内随时可用；沙箱外的任意目录需老板开「完全访问模式」）。",
        parameters={"type": "object", "properties": {
            "path": {"type": "string", "description": "要列的目录绝对路径（可选）；不传＝列当前工作目录，没设工作目录则列内容库"},
        }},
        handler=list_files,
        read_only=True,
    ),
    Tool(
        name="find_files",
        description="按文件名【递归找文件】（像 Claude Code 的 Glob）。给一个目录和文件名规则(glob)，"
                    "如在桌面找所有 *.xlsx、或找 **/采购* 这种。返回匹配到的文件完整路径。"
                    "沙箱内（内容库/老板选定的）随时可用；要在沙箱外任意目录找，需老板开「完全访问模式」。",
        parameters={"type": "object", "properties": {
            "root_path": {"type": "string", "description": "在哪个目录下找（绝对路径，如老板的桌面）"},
            "pattern": {"type": "string", "description": "文件名规则(glob)，如 *.xlsx、**/采购*、报表*.csv"},
            "max_results": {"type": "integer", "description": "最多返回几个（默认 100）"},
        }, "required": ["root_path", "pattern"]},
        handler=find_files,
        read_only=True,
    ),
    Tool(
        name="search_in_files",
        description="按【内容】搜文件（像 Claude Code 的 Grep）。给一个目录和关键词/正则，"
                    "搜出哪些文件里含这个词，返回 文件:行号:命中行。找『提到过XX的文件/报表里写了YY的』时用。"
                    "沙箱内随时可用；沙箱外任意目录需老板开「完全访问模式」。",
        parameters={"type": "object", "properties": {
            "root_path": {"type": "string", "description": "在哪个目录下搜（绝对路径）"},
            "query": {"type": "string", "description": "要搜的关键词或正则，如『提成』或『\\d{11}』"},
            "max_results": {"type": "integer", "description": "最多返回几条命中（默认 100）"},
        }, "required": ["root_path", "query"]},
        handler=search_in_files,
        read_only=True,
    ),
    Tool(
        name="run_command",
        description="在老板本机【跑一条命令】（像 Claude Code 的 Bash，最危险）。"
                    "只有老板开了「完全访问模式」才可用。逐项确认/自动接受档会把命令原文弹给老板确认才执行；"
                    "完全访问档则自己跑、不逐个问（命令会显示在对话里，老板可随时打断）。"
                    "禁止 && || | ; > < 反引号 这类拼接/重定向；危险命令(删根/提权/格式化等)任何档位都直接拒。"
                    "适合跑 ls/find/python 脚本/git status 这类一次性查看类命令。",
        parameters={"type": "object", "properties": {
            "command": {"type": "string", "description": "要执行的单条命令原文（不要用 && | ; 等拼接）"},
            "timeout_sec": {"type": "integer", "description": "超时秒数（默认 30，最多 300）"},
            "cwd": {"type": "string", "description": "工作目录（可选，绝对路径）"},
        }, "required": ["command"]},
        handler=run_command,
        requires_approval=True,
        approval_class="command",
        force_confirm=False,  # 跟权限档走:L1/L2弹卡确认、L3完全访问自己跑;危险黑名单在handler、与档位无关永远拦
        preview=preview_run_command,
        # 审批闸 2.0：①命中危险黑名单 → _plan_tool_call 直接拒绝、不占审批卡名额（判据同 handler 内部）；
        # ③命中安全前缀白名单 → 任何权限档都可自动放行、不弹卡。
        fatal_reason_for=_run_command_fatal_reason,
        safe_prefix_for=_run_command_safe_prefix,
    ),
    Tool(
        name="read_file",
        description="读取内容库里某个文件的内容（编辑前必须先读，才知道里面是什么）。Excel 列单元格；PDF/Word(.docx)/PPT(.pptx) 自动提取文字。",
        parameters={"type": "object", "properties": {"path": {"type": "string", "description": "内容库内的文件名/相对路径，或老板当场选定文件的完整路径"}}, "required": ["path"]},
        handler=read_file,
        read_only=True,
        requires_approval_for=lambda args, ctx: _is_sensitive_file(args.get("path")),
        approval_class="sensitive_read",
        approval_reason=_sensitive_read_reason,
        preview=_sensitive_read_preview,
    ),
    Tool(
        name="write_file",
        description="把内容写进本机内容库的一个文件（新建或覆盖，如保存一份文案/清单）。覆盖会先自动备份原件。",
        parameters={"type": "object", "properties": {"path": {"type": "string"}, "content": {"type": "string"}}, "required": ["path", "content"]},
        handler=write_file,
        requires_approval=True,
        approval_class="file",
        preview=preview_write_file,
    ),
    Tool(
        name="edit_file",
        description="修改一个文本文件的某一段：把指定原文精确替换成新文本。改前自动备份、可回滚。",
        parameters={"type": "object", "properties": {"path": {"type": "string"}, "old_text": {"type": "string"}, "new_text": {"type": "string"}}, "required": ["path", "old_text", "new_text"]},
        handler=edit_file,
        requires_approval=True,
        approval_class="file",
        preview=preview_edit_file,
    ),
    Tool(
        name="edit_excel",
        description="直接修改本机的 Excel 报表（改营业额、改某个数、加一列提成等）。先 read_file 看清单元格坐标，再给要改的单元格。改前自动备份、改后回传每格的前后对比。",
        parameters={"type": "object", "properties": {
            "path": {"type": "string"},
            "changes": {"type": "array", "items": {"type": "object", "properties": {
                "sheet": {"type": "string", "description": "工作表名，留空=第一个表"},
                "cell": {"type": "string", "description": "A1 式坐标，如 B2"},
                "value": {"description": "新值"},
            }, "required": ["cell", "value"]}},
        }, "required": ["path", "changes"]},
        handler=edit_excel,
        requires_approval=True,
        approval_class="file",
        preview=preview_edit_excel,
    ),
    Tool(
        name="delete_file",
        description="删除内容库里的一个文件（如清掉一份不要的旧文案/旧报表/旧清单）。删前自动备份、可恢复。"
                    "高危不可逆动作——【任何权限档都会先把「要删哪个文件」弹给老板确认，确认后才真正删】，"
                    "「跳过确认」档也不例外。只删单个文件、不删文件夹。",
        parameters={"type": "object", "properties": {
            "path": {"type": "string", "description": "要删的文件（内容库内文件名/相对路径，或老板当场选定文件的完整路径）"},
        }, "required": ["path"]},
        handler=delete_file,
        requires_approval=True,
        approval_class="file",
        force_confirm=True,  # 删除不可逆（虽自动备份）——任何档位都强制确认，不被「跳过确认」旁路
        preview=preview_delete_file,
    ),
]


def register_local_tools(registry=None) -> int:
    """把本地文件工具注册进注册表。仅桌面本地模式调用。返回注册数。"""
    reg = registry or default_registry
    for t in _LOCAL_TOOLS:
        if reg.get(t.name) is None:
            reg.register(t)
    return len(_LOCAL_TOOLS)


# 仅桌面全本地模式自动注册（云端 web 版不设 DESKTOP_LOCAL → 拿不到文件操作工具）
if os.environ.get("DESKTOP_LOCAL") == "1":
    register_local_tools()
    logger.info("已注册 %d 个本地文件操作工具（桌面全本地模式）", len(_LOCAL_TOOLS))
