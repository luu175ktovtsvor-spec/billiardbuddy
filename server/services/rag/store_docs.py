"""店铺资料库：老板选一个文件夹（合同/进货单/排班表/价目表…）→ 提取纯文本 → 分块 → 灌进
通用向量 RAG（source_type="store_doc"）→ 带出处检索。

跟 services/rag/recall.py 现有的 backfill_from_generations 是同一套底层(embedder/index_store)，
但这里是全新的一路输入(用户自己的文档，不是 AI 生成记录)，独立 source_type 隔离检索。

跟 services/agent/local_tools.py 的 read_file 不同：那条是"读给 Agent 当轮看"（有 9000 字截断 +
"--- 第N页 ---"这类给人看的装饰，直接拿去 embed 会丢内容+塞进噪声）；这里是"索引用途"的纯文本提取，
不截断、不装饰，长文档整篇分块全部索引。

分块策略：按自然单元（段落/PDF页/PPT页/Excel行组）先切，超长单元再硬切，最后用滑动缓冲区打包成
~400-800 字的块、块间留一点重叠（防在语义中间硬生生切断）。Excel 是例外——直接按"表头 + 若干行"
成块（带上表头才有语义，比逐格坐标转储管用）。
"""
import asyncio
import hashlib
import logging
from datetime import datetime, timezone
from pathlib import Path

# db.session.async_session：模块级引用，供测试 monkeypatch 成指向内存库的 sessionmaker
# （同 services/media_jobs_runner.py 的做法）——后台任务不能复用触发它的那个请求 session
# （请求早返回关闭了），run_folder_reindex_job 必须自己开一条。
from db.session import async_session
from services.rag import index_store
from services.rag.recall import index_text, recall

logger = logging.getLogger(__name__)

# 索引支持的文档格式（别引入新解析依赖：pypdf/python-docx/python-pptx/openpyxl 项目里已在用）。
SUPPORTED_EXTS = {".pdf", ".docx", ".pptx", ".xlsx", ".xlsm", ".txt", ".md"}

# 单文件读取护栏：太大的文件跳过（记 warning），防止老板选了个装了几百 GB 视频的文件夹时卡死。
_MAX_FILE_BYTES = 30 * 1024 * 1024  # 30MB

# 分块目标（字符数，中文场景 1 字≈1 token 量级，够用不用精确 tokenizer）。
_CHUNK_TARGET_MAX = 800
_CHUNK_OVERLAP = 80
# 打包完最后一块太小(< 这个数)且前面还有块时，并入上一块，防止出现"一两个字"的畸零块。
_CHUNK_MERGE_TAIL_BELOW = 150

# xlsx 每块打包的数据行数（表头随每块重复，保证每块单独看也有语义）。
_XLSX_ROWS_PER_CHUNK = 25


# ────────────────────────────── 纯文本提取（不截断/不加装饰，供索引用） ──────────────────────────────

def _extract_units_pdf(path: Path) -> list[str]:
    from pypdf import PdfReader
    reader = PdfReader(str(path))
    units = []
    for page in reader.pages:
        t = (page.extract_text() or "").strip()
        if t:
            units.append(t)
    return units


def _extract_units_docx(path: Path) -> list[str]:
    import docx
    doc = docx.Document(str(path))
    units = [p.text.strip() for p in doc.paragraphs if p.text.strip()]
    for tbl in doc.tables:
        for row in tbl.rows:
            cells = [c.text.strip() for c in row.cells]
            if any(cells):
                units.append(" | ".join(cells))
    return units


def _extract_units_pptx(path: Path) -> list[str]:
    from pptx import Presentation
    prs = Presentation(str(path))
    units = []
    for slide in prs.slides:
        texts = [s.text_frame.text.strip() for s in slide.shapes
                 if s.has_text_frame and s.text_frame.text.strip()]
        if texts:
            units.append("\n".join(texts))
    return units


def _extract_units_text(path: Path) -> list[str]:
    try:
        text = path.read_text(encoding="utf-8")
    except UnicodeDecodeError:
        text = path.read_text(encoding="utf-8", errors="ignore")
    # 按空行分段；没有空行分隔就整篇当一个单元（超长交给 _hard_split 兜底，不会丢尾部）。
    import re
    paras = [p.strip() for p in re.split(r"\n\s*\n", text) if p.strip()]
    if paras:
        return paras
    return [text.strip()] if text.strip() else []


def _extract_chunks_xlsx(path: Path) -> list[str]:
    """xlsx 专属：不走通用分块器，直接按"表头 + 若干行"成块——带表头让每块脱离原表也能看懂，
    比给人看的坐标转储(A1=xxx)更有语义、也更利于 embedding 命中。"""
    from openpyxl import load_workbook
    wb = load_workbook(str(path), data_only=True)
    chunks: list[str] = []
    for ws in wb.worksheets:
        rows_iter = ws.iter_rows(values_only=True)
        try:
            header_row = next(rows_iter)
        except StopIteration:
            continue
        header_cells = [str(c).strip() if c is not None else "" for c in header_row]
        header_line = f"工作表「{ws.title}」字段：" + "、".join(h for h in header_cells if h)

        batch: list[str] = []
        for row in rows_iter:
            if row is None or all(c is None or str(c).strip() == "" for c in row):
                continue
            cells = []
            for h, v in zip(header_cells, row):
                if v is None or str(v).strip() == "":
                    continue
                cells.append(f"{h}={v}" if h else str(v))
            if not cells:
                continue
            batch.append("；".join(cells))
            if len(batch) >= _XLSX_ROWS_PER_CHUNK:
                chunks.append(header_line + "\n" + "\n".join(batch))
                batch = []
        if batch:
            chunks.append(header_line + "\n" + "\n".join(batch))
    return chunks


def _hard_split(text: str, size: int, overlap: int) -> list[str]:
    """把超长文本按固定长度切片，切片间保留 overlap 字符重叠（防切断语义、防丢内容）。"""
    if len(text) <= size:
        return [text]
    step = max(size - overlap, 1)
    out = []
    i = 0
    n = len(text)
    while i < n:
        out.append(text[i:i + size])
        if i + size >= n:
            break
        i += step
    return out


def _pack_units_into_chunks(units: list[str], target_max: int = _CHUNK_TARGET_MAX,
                             overlap: int = _CHUNK_OVERLAP) -> list[str]:
    """把一串"自然单元"（段落/页/表格行…）打包成目标大小的块，块间留重叠。
    先把超长单元硬切成 <=target_max 的片段，再用滑动缓冲区顺序打包——保证全篇内容都被覆盖，
    最后剩的缓冲区一定会被 flush 成最后一块（不丢尾部）。"""
    pieces: list[str] = []
    for u in units:
        u = (u or "").strip()
        if not u:
            continue
        if len(u) > target_max:
            pieces.extend(_hard_split(u, target_max, overlap))
        else:
            pieces.append(u)

    chunks: list[str] = []
    buf = ""
    for piece in pieces:
        if not buf:
            buf = piece
            continue
        if len(buf) + 1 + len(piece) <= target_max:
            buf = buf + "\n" + piece
        else:
            chunks.append(buf)
            tail = buf[-overlap:] if overlap and len(buf) > overlap else buf
            buf = tail + "\n" + piece
    if buf.strip():
        chunks.append(buf)

    # 畸零尾块并入前一块（纯体验优化，不影响"不丢内容"——内容仍在，只是不单独成块）。
    if len(chunks) > 1 and len(chunks[-1]) < _CHUNK_MERGE_TAIL_BELOW:
        chunks[-2] = chunks[-2] + "\n" + chunks[-1]
        chunks.pop()
    return chunks


def extract_and_chunk(path) -> list[str]:
    """纯文本提取(不截断/不加装饰) + 语义分块。按后缀分派；不支持的后缀返回空列表。

    xlsx 直接返回"表头+行组"块（不走通用打包器）；其余格式先按自然单元提取、再打包成块。
    调用方（index_store_docs_folder）负责 try/except 兜底单文件解析失败，这里让异常自然抛出。
    """
    p = Path(path)
    suffix = p.suffix.lower()
    if suffix in (".xlsx", ".xlsm"):
        return _extract_chunks_xlsx(p)
    if suffix == ".pdf":
        units = _extract_units_pdf(p)
    elif suffix == ".docx":
        units = _extract_units_docx(p)
    elif suffix == ".pptx":
        units = _extract_units_pptx(p)
    elif suffix in (".txt", ".md"):
        units = _extract_units_text(p)
    else:
        return []
    return _pack_units_into_chunks(units)


# ────────────────────────────── 索引管线：遍历文件夹 → 分块 → 灌进向量库 ──────────────────────────────

def _parse_source_id(source_id: str) -> tuple[str, int]:
    """source_id 编码为 "<相对路径>#<chunk序号>"——用 rsplit 只切最后一个 '#'，防路径本身含 '#'。"""
    rel, _, idx = source_id.rpartition("#")
    try:
        return rel, int(idx)
    except (ValueError, TypeError):
        return source_id, 0


def _file_fingerprint(path: Path) -> str:
    st = path.stat()
    return f"{st.st_mtime_ns}:{st.st_size}"


def _iter_supported_files(folder: Path):
    for p in sorted(folder.rglob("*")):
        if not p.is_file():
            continue
        if p.suffix.lower() not in SUPPORTED_EXTS:
            continue
        if any(part.startswith(".") for part in p.relative_to(folder).parts):
            continue  # 跳过隐藏文件/隐藏目录(.git/.DS_Store 等)
        yield p


def index_store_docs_folder(store_id, folder_path: str, progress_cb=None) -> dict:
    """遍历文件夹索引店铺资料。返回 {file_count, chunk_count, errors, fatal_error}：
    - file_count/chunk_count：跑完后这个 store 在 store_doc 索引里【当前总共】有多少文件/多少块
      （覆盖新增+未变+减去被清掉的陈旧块，是"现状"而非"本次新增了多少"）。
    - errors：单个文件解析失败的提示(不中断整批，跳过继续)。
    - fatal_error：文件夹本身就打不开/不是目录——整批没法跑，调用方据此把状态标 error。

    增量：同一文件同一块内容不变（file 的 mtime+size 与 chunk 内容都一致）就跳过重嵌，
    不用每次全量重算 embedding。文件改小/删除后，旧的多余 chunk 会被清掉，不留陈旧搜索结果。
    故障安全：单个文件解析失败只记 warning、跳过，不中断整批索引。
    """
    folder = Path(folder_path)
    if not folder.exists() or not folder.is_dir():
        return {"file_count": 0, "chunk_count": 0, "errors": [],
                "fatal_error": f"文件夹不存在或不可读：{folder_path}"}

    files = list(_iter_supported_files(folder))
    total = len(files)
    errors: list[str] = []
    existing_fps = index_store.existing_fingerprints(str(store_id), "store_doc")  # 一次性取，增量比对用
    seen_rel_paths: set[str] = set()

    for idx, path in enumerate(files, 1):
        rel = path.relative_to(folder).as_posix()
        seen_rel_paths.add(rel)
        try:
            if path.stat().st_size > _MAX_FILE_BYTES:
                errors.append(f"{rel}：文件过大(>{_MAX_FILE_BYTES // (1024 * 1024)}MB)，已跳过")
                continue
            file_fp = _file_fingerprint(path)
        except OSError as e:
            errors.append(f"{rel}：读取文件信息失败（{e}）")
            continue

        try:
            chunks = extract_and_chunk(path)
        except Exception as e:  # noqa: BLE001 —— 故障安全：单文件解析失败不能中断整批索引
            logger.warning("店铺资料索引：文件解析失败，跳过 %s", rel, exc_info=True)
            errors.append(f"{rel}：解析失败（{e}）")
            continue

        new_ids = set()
        for i, chunk in enumerate(chunks):
            source_id = f"{rel}#{i}"
            new_ids.add(source_id)
            fp = hashlib.sha1(f"{file_fp}|{i}|{chunk}".encode("utf-8")).hexdigest()
            if existing_fps.get(source_id) == fp:
                continue  # 内容没变，跳过重嵌
            meta = {"file_name": path.name, "file_path": rel, "chunk_index": i}
            index_text(store_id, "store_doc", source_id, chunk, ts="", meta=meta, fp=fp)

        # 清掉这个文件里"上次有、这次没了"的陈旧块（文件被改短/段落被删）。
        stale = [sid for sid in existing_fps if sid.rpartition("#")[0] == rel and sid not in new_ids]
        if stale:
            index_store.delete_ids(str(store_id), "store_doc", stale)

        if progress_cb:
            try:
                progress_cb(idx, total, rel)
            except Exception:  # noqa: BLE001 —— 进度回调是尽力而为，不能拖累索引本身
                pass

    # 清掉"文件夹里已经没有这个文件了"的陈旧块(老板删了源文件)。
    orphan_files = [sid for sid in existing_fps if sid.rpartition("#")[0] not in seen_rel_paths]
    if orphan_files:
        index_store.delete_ids(str(store_id), "store_doc", orphan_files)

    final_fps = index_store.existing_fingerprints(str(store_id), "store_doc")
    chunk_count = len(final_fps)
    file_count = len({sid.rpartition("#")[0] for sid in final_fps})
    return {"file_count": file_count, "chunk_count": chunk_count, "errors": errors, "fatal_error": None}


# ────────────────────────────── 检索：只召回 store_doc，带出处 ──────────────────────────────

def search_store_docs_impl(store_id, query: str, top: int = 5) -> list[dict]:
    """检索店铺资料（只召回 source_type="store_doc"，绝不把老板过去的生成记录也搜出来）。
    返回 [{file_name, file_path, chunk_index, snippet, score}]，可直接引用出处。"""
    if not query or not query.strip():
        return []
    hits = recall(str(store_id), query, top=top, source_type="store_doc")
    results = []
    for h in hits:
        rel_path, chunk_idx = _parse_source_id(h["source_id"])
        snippet = " ".join((h["text"] or "").split())[:280]
        results.append({
            "file_name": Path(rel_path).name,
            "file_path": rel_path,
            "chunk_index": chunk_idx,
            "snippet": snippet,
            "score": h["score"],
        })
    return results


# ────────────────────────────── 后台任务：索引大文件夹可能慢，交后台跑 + 写回配置表 ──────────────────────────────

async def run_folder_reindex_job(store_id, folder_path: str) -> None:
    """后台任务体：真正跑索引 + 把结果/状态写回 StoreDocLibrary 配置表。
    自己开 async_session（不能复用触发它的那个请求 session，请求早返回了）。"""
    from models.store_doc_library import StoreDocLibrary
    from sqlalchemy import select

    try:
        stats = await asyncio.to_thread(index_store_docs_folder, store_id, folder_path)
    except Exception as e:  # noqa: BLE001 —— 索引管线本身炸了也要把 error 落回配置表，不能崩后台任务、更不能吞掉
        logger.exception("店铺资料索引后台任务失败 store=%s folder=%s", store_id, folder_path)
        stats = {"file_count": 0, "chunk_count": 0, "errors": [], "fatal_error": str(e) or e.__class__.__name__}

    fatal = stats.get("fatal_error")
    if fatal:
        status, last_error = "error", fatal
    else:
        errors = stats.get("errors") or []
        status = "ready"
        last_error = "；".join(errors)[:2000] if errors else None

    async with async_session() as db:
        row = (await db.execute(
            select(StoreDocLibrary).where(StoreDocLibrary.store_id == store_id)
        )).scalars().first()
        if row is None:
            return  # 配置行被老板中途清掉了——静默放弃写回，不重建
        row.status = status
        row.last_error = last_error
        if not fatal:
            row.indexed_file_count = stats.get("file_count", 0)
            row.indexed_chunk_count = stats.get("chunk_count", 0)
            row.last_indexed_at = datetime.now(timezone.utc)
        row.updated_at = datetime.now(timezone.utc)
        await db.commit()
