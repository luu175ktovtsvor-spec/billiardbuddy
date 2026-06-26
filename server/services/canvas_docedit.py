"""Word(.docx) / PPT(.pptx) 「文字级编辑 + 写回原文件」。

思路（关键）：把文档拆成带【稳定 id】的文本块给前端编辑；保存时按 id 回到原文件、
【原地只改这几段的文字】、保留文档其余结构与格式，再写回（改前自动备份）。
不重造文档 = 真·写回，符合"只改我改的、别动其它"。

扩展：要支持别的可按块编辑的格式，加一对 _xxx_blocks / _xxx_write 即可。
"""
from __future__ import annotations

import re
from pathlib import Path


# ───────────────── 读：文档 → 带 id 的文本块 ─────────────────

def _para_kind(style_name: str) -> str:
    s = (style_name or "").lower()
    if s.startswith("heading 1") or s == "title":
        return "h1"
    if s.startswith("heading 2"):
        return "h2"
    if s.startswith("heading 3"):
        return "h3"
    if "list" in s or s.startswith("bullet"):
        return "li"
    return "p"


def _docx_blocks(path: Path) -> list[dict]:
    import docx

    d = docx.Document(str(path))
    blocks: list[dict] = []
    for i, p in enumerate(d.paragraphs):
        if not (p.text or "").strip():
            continue
        blocks.append({"id": f"p{i}", "kind": _para_kind(p.style.name if p.style else ""), "text": p.text})
    for ti, tbl in enumerate(d.tables):
        for ri, row in enumerate(tbl.rows):
            for ci, cell in enumerate(row.cells):
                if (cell.text or "").strip():
                    blocks.append({"id": f"t{ti}-{ri}-{ci}", "kind": "cell", "text": cell.text})
    return blocks


def _pptx_blocks(path: Path) -> list[dict]:
    from pptx import Presentation

    prs = Presentation(str(path))
    blocks: list[dict] = []
    for si, slide in enumerate(prs.slides):
        title_id = None
        try:
            if slide.shapes.title is not None:
                title_id = slide.shapes.title.shape_id
        except Exception:  # noqa: BLE001
            title_id = None
        for shi, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            is_title = title_id is not None and shape.shape_id == title_id
            for pi, para in enumerate(shape.text_frame.paragraphs):
                if not (para.text or "").strip():
                    continue
                blocks.append({
                    "id": f"s{si}-{shi}-{pi}",
                    "kind": "title" if is_title else "body",
                    "text": para.text,
                    "slide": si + 1,
                })
    return blocks


def read_blocks(path: Path) -> dict:
    """读 docx/pptx → {kind, blocks:[{id, kind, text, slide?}]}。"""
    ext = path.suffix.lower()
    if ext == ".docx":
        return {"kind": "docx", "blocks": _docx_blocks(path)}
    if ext == ".pptx":
        return {"kind": "pptx", "blocks": _pptx_blocks(path)}
    raise ValueError(f"不支持按块编辑：{ext}")


# ───────────────── 写：把改动按 id 原地写回原文件 ─────────────────

def _set_para_runs(para, text: str) -> None:
    """把段落文字换成 text，尽量保留每个 run 原有的内联格式（加粗/颜色/字号）和超链接。
    算法：找新旧文本的公共前缀 + 后缀，只动变化区间里的 run，未变区间的 run 原封不动。
    同时兼容 docx（w:r + w:hyperlink）和 pptx（a:r）段落。"""
    # 判断 docx 还是 pptx（XML 命名空间不同）
    tag = para._element.tag or ""
    is_docx = "wordprocessingml" in tag

    if is_docx:
        from docx.oxml.ns import qn
        from docx.text.run import Run as _Run
        all_runs = []
        for child in para._element:
            if child.tag == qn("w:r"):
                all_runs.append(_Run(child, para))
            elif child.tag == qn("w:hyperlink"):
                for r_elem in child.findall(qn("w:r")):
                    all_runs.append(_Run(r_elem, para))
    else:
        all_runs = list(para.runs)

    if not all_runs:
        if is_docx:
            para.add_run(text)
        else:
            r = para.add_run()
            r.text = text
        return

    old_text = "".join(r.text or "" for r in all_runs)
    if old_text == text:
        return

    # 公共前缀长度
    plen = 0
    for a, b in zip(old_text, text):
        if a != b:
            break
        plen += 1

    # 公共后缀长度（前缀之后的部分）
    slen = 0
    if plen < len(old_text) and plen < len(text):
        for a, b in zip(reversed(old_text[plen:]), reversed(text[plen:])):
            if a != b:
                break
            slen += 1

    new_mid = text[plen : len(text) - slen] if slen else text[plen:]
    old_cs = plen
    old_ce = len(old_text) - slen

    # 每个 run 对应 [start, end) 字符位置
    pos = 0
    run_spans = []
    for r in all_runs:
        rlen = len(r.text or "")
        run_spans.append((pos, pos + rlen, r))
        pos += rlen

    placed = False
    if old_cs == old_ce:
        # 纯插入：找到插入点所在的 run，把新文本嵌入
        for start, end, run in run_spans:
            if start <= old_cs <= end:
                offset = old_cs - start
                rtxt = run.text or ""
                run.text = rtxt[:offset] + new_mid + rtxt[offset:]
                placed = True
                break
    else:
        for start, end, run in run_spans:
            if end <= old_cs or start >= old_ce:
                continue
            rtxt = run.text or ""
            before = rtxt[: max(0, old_cs - start)]
            after = rtxt[min(len(rtxt), old_ce - start) :]
            if not placed:
                run.text = before + new_mid + after
                placed = True
            else:
                run.text = after

    if not placed:
        all_runs[0].text = text
        for r in all_runs[1:]:
            r.text = ""


def _docx_write(path: Path, edits: dict) -> None:
    import docx

    d = docx.Document(str(path))
    paras = d.paragraphs
    for key, new in edits.items():
        if key.startswith("p"):
            i = int(key[1:])
            if 0 <= i < len(paras):
                _set_para_runs(paras[i], new)
        elif key.startswith("t"):
            m = re.match(r"t(\d+)-(\d+)-(\d+)$", key)
            if not m:
                continue
            ti, ri, ci = (int(x) for x in m.groups())
            try:
                cell = d.tables[ti].rows[ri].cells[ci]
            except (IndexError, AttributeError):
                continue
            # 单元格：改第一个段落，清掉其余段落文字
            cps = cell.paragraphs
            if cps:
                _set_para_runs(cps[0], new)
                for extra in cps[1:]:
                    for r in extra.runs:
                        r.text = ""
    d.save(str(path))


def _pptx_write(path: Path, edits: dict) -> None:
    from pptx import Presentation

    prs = Presentation(str(path))
    slides = list(prs.slides)
    for key, new in edits.items():
        m = re.match(r"s(\d+)-(\d+)-(\d+)$", key)
        if not m:
            continue
        si, shi, pi = (int(x) for x in m.groups())
        if si >= len(slides):
            continue
        shapes = list(slides[si].shapes)
        if shi >= len(shapes) or not shapes[shi].has_text_frame:
            continue
        paras = shapes[shi].text_frame.paragraphs
        if pi >= len(paras):
            continue
        _set_para_runs(paras[pi], new)
    prs.save(str(path))


def write_blocks(path: Path, edits: dict) -> None:
    """把 {id: 新文字} 原地写回原文件（改前自动备份）。"""
    from services.agent.local_tools import _backup

    if not edits:
        return
    ext = path.suffix.lower()
    if ext not in (".docx", ".pptx"):
        raise ValueError(f"不支持写回：{ext}")
    _backup(path)
    if ext == ".docx":
        _docx_write(path, edits)
    else:
        _pptx_write(path, edits)
