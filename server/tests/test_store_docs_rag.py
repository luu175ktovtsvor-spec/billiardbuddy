"""D-Task-5 店铺资料库：分块索引 + 带出处检索。

锁住：
- extract_and_chunk：pdf/docx/pptx/xlsx/txt 各能提取+分块；长文档不因套用旧的 9000 字截断而丢尾部内容；
  xlsx 按"表头+若干行"成块（带表头有语义，不是坐标转储）。
- index_store_docs_folder：索引一个临时文件夹 → vectors 表有 store_doc 记录、file/chunk 计数对；
  增量（同内容第二次跑不重嵌）；单文件解析失败不中断整批；文件被删/内容变短后陈旧 chunk 被清掉。
- search_store_docs_impl：只召回 store_doc（不管 generation），带出处(file_name)。
- source_type 双向隔离：search_store_docs 不串 generation；recall_my_content(source_type="generation")
  也不该被 store_doc 污染（这条是本单排查出的必然连带修复，不只是"新增"单向隔离）。
- 跨店隔离：A 店索引的文档，B 店搜不到。
- 工具元信息：search_store_docs 登记在 default_registry、read_only、无 approval、concurrent_safe、
  不在 BILLIARDS_TOOL_NAMES。

用确定性词面嵌入器(DeterministicEmbedder)+ 独立 tmp 向量库，稳定可重复，同 test_rag.py / test_m11_rag_fixes.py 的约定。
"""
import pytest

from services.rag.store_docs import (
    extract_and_chunk,
    index_store_docs_folder,
    search_store_docs_impl,
)


@pytest.fixture(autouse=True)
def _isolated_index(tmp_path, monkeypatch):
    """向量库指向 tmp、用确定性词面后端、清连接缓存，隔离真实环境与跨测试污染。"""
    monkeypatch.setenv("DESKTOP_RAG_DIR", str(tmp_path / "rag"))
    monkeypatch.delenv("RAG_EMBEDDER", raising=False)
    from services.rag import index_store
    index_store.reset_for_test()
    from services.rag.embedder import DeterministicEmbedder
    import services.rag.embedder as emb
    emb._embedder = DeterministicEmbedder()
    yield
    index_store.reset_for_test()
    emb._embedder = DeterministicEmbedder()


# ══════════════════════════════ extract_and_chunk：各格式提取 + 分块 ══════════════════════════════

def test_extract_and_chunk_txt_basic_paragraphs(tmp_path):
    p = tmp_path / "笔记.txt"
    p.write_text("第一段内容说明。\n\n第二段内容说明。\n\n第三段内容说明。", encoding="utf-8")
    chunks = extract_and_chunk(p)
    joined = "\n".join(chunks)
    assert "第一段内容说明" in joined
    assert "第二段内容说明" in joined
    assert "第三段内容说明" in joined


def test_extract_and_chunk_long_txt_does_not_lose_tail_content(tmp_path):
    """长文档必须整篇分块全部索引——不能像 read_file 那样套 9000 字硬截断丢掉后面内容。"""
    paras = [f"第{i}段：" + ("这是用来填充长度的正文内容。" * 60) for i in range(30)]
    paras[-1] += " TAIL_MARKER_UNIQUE_9F3A"
    text = "\n\n".join(paras)
    assert len(text) > 9000, "样本必须比旧 _DOC_MAX_CHARS(9000) 更长，才能验证不会被那套截断坑"

    p = tmp_path / "长笔记.txt"
    p.write_text(text, encoding="utf-8")
    chunks = extract_and_chunk(p)

    assert len(chunks) > 1, "长文档应该被切成多块"
    assert any("TAIL_MARKER_UNIQUE_9F3A" in c for c in chunks), "长文档分块后尾部内容不应丢失"
    assert any("第0段" in c for c in chunks), "开头内容也该在"
    # 分块总长度不该比原文短太多（允许因重叠而略长，绝不允许因截断而大幅变短）
    assert sum(len(c) for c in chunks) >= len(text) * 0.9


def test_extract_and_chunk_docx_paragraphs_and_tables(tmp_path):
    import docx
    d = docx.Document()
    d.add_paragraph("合同第一条：租期一年。")
    d.add_paragraph("合同第二条：押金两千元。")
    t = d.add_table(rows=1, cols=2)
    t.rows[0].cells[0].text = "项目"
    t.rows[0].cells[1].text = "金额"
    p = tmp_path / "合同.docx"
    d.save(str(p))

    chunks = extract_and_chunk(p)
    joined = "\n".join(chunks)
    assert "租期一年" in joined
    assert "押金两千元" in joined
    assert "项目 | 金额" in joined


def test_extract_and_chunk_pptx_slides(tmp_path):
    from pptx import Presentation
    prs = Presentation()
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "排班表说明"
    body = slide.placeholders[1].text_frame
    body.text = "早班9点到17点"
    body.add_paragraph().text = "晚班17点到1点"
    p = tmp_path / "排班.pptx"
    prs.save(str(p))

    chunks = extract_and_chunk(p)
    joined = "\n".join(chunks)
    assert "排班表说明" in joined
    assert "早班9点到17点" in joined
    assert "晚班17点到1点" in joined


def test_extract_and_chunk_xlsx_headers_and_rows(tmp_path):
    """xlsx 按"表头+若干行"成块——带表头才有语义，不是给人看的坐标转储(A1=xxx)。"""
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "价目表"
    ws.append(["项目", "价格"])
    for i in range(3):
        ws.append([f"项目{i}", 100 + i])
    p = tmp_path / "价目表.xlsx"
    wb.save(str(p))

    chunks = extract_and_chunk(p)
    assert len(chunks) == 1
    assert "价目表" in chunks[0]
    assert "项目" in chunks[0] and "价格" in chunks[0]
    assert "项目0" in chunks[0] and "100" in chunks[0]


def test_extract_and_chunk_xlsx_batches_many_rows_with_repeated_header(tmp_path):
    from openpyxl import Workbook
    wb = Workbook()
    ws = wb.active
    ws.title = "排班表"
    ws.append(["日期", "班次"])
    for i in range(60):
        ws.append([f"7-{i}", "早班"])
    p = tmp_path / "排班表.xlsx"
    wb.save(str(p))

    chunks = extract_and_chunk(p)
    assert len(chunks) == 3  # 25 + 25 + 10
    for c in chunks:
        assert "字段：" in c, "每块都该带表头，脱离原表也能看懂"


def test_extract_and_chunk_pdf_multi_page(monkeypatch, tmp_path):
    """PDF 生成没有零依赖库(pypdf 本身不写文字内容)，monkeypatch pypdf.PdfReader 隔离测试
    我们自己的分页提取+分块逻辑（pypdf.extract_text() 本身的正确性不在本单测试范围）。"""
    import pypdf

    class _FakePage:
        def __init__(self, text):
            self._text = text

        def extract_text(self):
            return self._text

    class _FakeReader:
        def __init__(self, _path):
            self.pages = [_FakePage("第一页：合同背景介绍。"), _FakePage("第二页：违约责任条款。")]

    monkeypatch.setattr(pypdf, "PdfReader", _FakeReader)
    p = tmp_path / "合同.pdf"
    p.write_bytes(b"%PDF-fake-content-not-real")

    chunks = extract_and_chunk(p)
    joined = "\n".join(chunks)
    assert "合同背景介绍" in joined
    assert "违约责任条款" in joined


def test_extract_and_chunk_unsupported_ext_returns_empty(tmp_path):
    p = tmp_path / "图片.png"
    p.write_bytes(b"\x89PNG")
    assert extract_and_chunk(p) == []


# ══════════════════════════════ index_store_docs_folder：索引管线 ══════════════════════════════

def test_index_store_docs_folder_indexes_files_and_counts(tmp_path):
    folder = tmp_path / "资料"
    folder.mkdir()
    (folder / "价目表.txt").write_text("单人台球 30 元一小时\n\n双人台球 50 元一小时", encoding="utf-8")
    (folder / "排班.txt").write_text("早班 9 点到 17 点\n\n晚班 17 点到 1 点", encoding="utf-8")

    sid = "store-basic"
    stats = index_store_docs_folder(sid, str(folder))
    assert stats["fatal_error"] is None
    assert stats["file_count"] == 2
    assert stats["chunk_count"] >= 2

    from services.rag import index_store
    rows = index_store._conn().execute(
        "SELECT source_id FROM vectors WHERE store_id=? AND source_type='store_doc'", (sid,)
    ).fetchall()
    assert len(rows) == stats["chunk_count"]


def test_index_store_docs_folder_missing_folder_is_fatal():
    stats = index_store_docs_folder("store-x", "/definitely/not/a/real/folder/xyz")
    assert stats["fatal_error"] is not None
    assert stats["file_count"] == 0 and stats["chunk_count"] == 0


def test_index_store_docs_folder_skips_unchanged_on_second_run(tmp_path, monkeypatch):
    folder = tmp_path / "资料"
    folder.mkdir()
    (folder / "价目.txt").write_text("单人 30 元一小时", encoding="utf-8")
    sid = "store-inc"

    stats1 = index_store_docs_folder(sid, str(folder))
    assert stats1["chunk_count"] == 1

    calls = {"n": 0}
    from services.rag.embedder import DeterministicEmbedder
    real_embed = DeterministicEmbedder.embed

    def _counting_embed(self, text):
        calls["n"] += 1
        return real_embed(self, text)

    monkeypatch.setattr(DeterministicEmbedder, "embed", _counting_embed)

    stats2 = index_store_docs_folder(sid, str(folder))
    assert calls["n"] == 0, "文件内容没变，第二次跑不该重新计算 embedding"
    assert stats2["chunk_count"] == 1


def test_index_store_docs_folder_reembeds_when_file_changes(tmp_path, monkeypatch):
    folder = tmp_path / "资料"
    folder.mkdir()
    f = folder / "价目.txt"
    f.write_text("单人 30 元一小时", encoding="utf-8")
    sid = "store-changed"
    index_store_docs_folder(sid, str(folder))

    import time
    time.sleep(0.01)
    f.write_text("单人 35 元一小时（涨价了）", encoding="utf-8")

    calls = {"n": 0}
    from services.rag.embedder import DeterministicEmbedder
    real_embed = DeterministicEmbedder.embed

    def _counting_embed(self, text):
        calls["n"] += 1
        return real_embed(self, text)

    monkeypatch.setattr(DeterministicEmbedder, "embed", _counting_embed)
    index_store_docs_folder(sid, str(folder))
    assert calls["n"] >= 1, "文件内容变了应该重新嵌入"


def test_index_store_docs_folder_skips_broken_file_without_aborting_batch(tmp_path):
    folder = tmp_path / "资料"
    folder.mkdir()
    (folder / "好的.txt").write_text("这是可以正常索引的内容", encoding="utf-8")
    (folder / "坏的.pdf").write_bytes(b"not a real pdf at all")

    sid = "store-fault"
    stats = index_store_docs_folder(sid, str(folder))
    assert stats["fatal_error"] is None, "单文件解析失败不该让整批失败"
    assert stats["file_count"] == 1, "只有'好的.txt'真正索引成功"
    assert any("坏的.pdf" in e for e in stats["errors"])


def test_index_store_docs_folder_prunes_chunks_of_deleted_file(tmp_path):
    folder = tmp_path / "资料"
    folder.mkdir()
    f = folder / "旧文件.txt"
    f.write_text("这是旧文件的内容，以后会被删掉", encoding="utf-8")
    sid = "store-prune"

    index_store_docs_folder(sid, str(folder))
    from services.rag import index_store
    before = index_store.existing_fingerprints(sid, "store_doc")
    assert any(k.startswith("旧文件.txt#") for k in before)

    f.unlink()
    index_store_docs_folder(sid, str(folder))
    after = index_store.existing_fingerprints(sid, "store_doc")
    assert not any(k.startswith("旧文件.txt#") for k in after), "文件被删后，陈旧 chunk 不该继续留在索引里"


def test_index_store_docs_folder_prunes_stale_chunks_when_file_shrinks(tmp_path):
    """文件从"很多段"改成"很少段"后，多出来的旧 chunk 该被清掉，否则会一直命中已经不存在的旧内容。"""
    folder = tmp_path / "资料"
    folder.mkdir()
    f = folder / "笔记.txt"
    # 每段填充到 ~700 字（够大才会各自独立成块，不会被打包器合并进同一块）
    paras = [f"第{i}段：" + ("独立" * 350) for i in range(10)]
    f.write_text("\n\n".join(paras), encoding="utf-8")
    sid = "store-shrink"
    stats1 = index_store_docs_folder(sid, str(folder))
    assert stats1["chunk_count"] >= 5

    import time
    time.sleep(0.01)
    f.write_text("只剩一小段内容了", encoding="utf-8")
    stats2 = index_store_docs_folder(sid, str(folder))
    assert stats2["chunk_count"] == 1

    from services.rag import index_store
    remaining = index_store.existing_fingerprints(sid, "store_doc")
    assert len(remaining) == 1


# ══════════════════════════════ search_store_docs_impl：带出处检索 + source_type 隔离 ══════════════════════════════

def test_search_store_docs_impl_returns_relevant_with_file_name(tmp_path):
    folder = tmp_path / "资料"
    folder.mkdir()
    (folder / "价目表.txt").write_text("单人台球 30 元一小时，双人台球 50 元一小时", encoding="utf-8")
    (folder / "排班表.txt").write_text("助教早班 9 点到 17 点上岗", encoding="utf-8")
    sid = "store-search"
    index_store_docs_folder(sid, str(folder))

    hits = search_store_docs_impl(sid, "台球一小时多少钱", top=3)
    assert hits
    assert hits[0]["file_name"] == "价目表.txt"
    assert "元" in hits[0]["snippet"]
    assert hits[0]["file_path"] == "价目表.txt"
    assert isinstance(hits[0]["chunk_index"], int)


def test_search_store_docs_impl_only_returns_store_doc_not_generation(tmp_path):
    """先塞一条 generation，再索引店铺资料，search_store_docs 不该把 generation 内容搜出来。"""
    from services.rag.recall import index_text
    sid = "store-isolation"
    index_text(sid, "generation", "g1", "双十一活动的历史生成文案，全场五折")

    folder = tmp_path / "资料"
    folder.mkdir()
    (folder / "活动细则.txt").write_text("双十一活动合同细则：满一百减三十", encoding="utf-8")
    index_store_docs_folder(sid, str(folder))

    hits = search_store_docs_impl(sid, "双十一活动", top=5)
    assert hits
    assert all(h["file_name"] == "活动细则.txt" for h in hits), "不该混进 generation 的内容"


def test_recall_my_content_source_type_does_not_leak_store_doc(tmp_path):
    """反过来也要隔离：recall_my_content(翻旧生成记录) 不该被店铺资料(store_doc)污染。
    这是本单必然的连带修复——index_store.search 原本不分 source_type，加了 store_doc 这第二路
    来源后，如果不显式过滤，老路径(recall_my_content)会开始搜到混进来的文档内容。"""
    from services.rag.recall import index_text, recall
    sid = "store-reverse-isolation"
    index_text(sid, "generation", "g1", "双十一活动历史文案，全场五折")

    folder = tmp_path / "资料"
    folder.mkdir()
    (folder / "合同.txt").write_text("双十一活动合同条款细则", encoding="utf-8")
    index_store_docs_folder(sid, str(folder))

    hits = recall(sid, "双十一活动", top=5, source_type="generation")
    assert hits
    assert all(h["source_type"] == "generation" for h in hits), "recall_my_content 不该被 store_doc 内容污染"


def test_cross_store_isolation_for_store_docs(tmp_path):
    folder_a = tmp_path / "A资料"
    folder_a.mkdir()
    (folder_a / "甲店合同.txt").write_text("甲店场地租期两年", encoding="utf-8")
    index_store_docs_folder("store-A", str(folder_a))

    hits_b = search_store_docs_impl("store-B", "场地租期", top=5)
    assert hits_b == [], "A 店索引的文档，B 店不该搜到"

    hits_a = search_store_docs_impl("store-A", "场地租期", top=5)
    assert hits_a and hits_a[0]["file_name"] == "甲店合同.txt"


def test_search_store_docs_impl_empty_query_returns_empty(tmp_path):
    assert search_store_docs_impl("store-empty", "", top=5) == []
    assert search_store_docs_impl("store-empty", "   ", top=5) == []


# ══════════════════════════════ 工具元信息：search_store_docs ══════════════════════════════

def test_search_store_docs_tool_metadata():
    from services.agent.local_tools import register_local_tools
    from services.agent.registry import BILLIARDS_TOOL_NAMES, ToolRegistry

    reg = ToolRegistry()
    register_local_tools(reg)
    t = reg.get("search_store_docs")
    assert t is not None, "search_store_docs 应登记进本地工具表"
    assert t.read_only is True
    assert t.requires_approval is False
    assert t.concurrent_safe is True
    assert "search_store_docs" not in BILLIARDS_TOOL_NAMES, "不是台球专属工具，通用模式也该能用"
