"""Word/PPT 按块读 + 写回 单测：读出块 → 改其中一块 → 写回 → 复读验证只改了那块。"""
import docx
from pptx import Presentation

from services import canvas_docedit


def test_docx_read_edit_writeback(tmp_path, monkeypatch):
    monkeypatch.setenv("DESKTOP_LIBRARY_DIR", str(tmp_path / "lib"))
    p = tmp_path / "方案.docx"
    d = docx.Document()
    d.add_heading("活动方案", level=1)
    d.add_paragraph("周末满 100 送 30。")
    d.add_paragraph("仅限堂食。")
    d.save(str(p))

    data = canvas_docedit.read_blocks(p)
    assert data["kind"] == "docx"
    texts = {b["id"]: b["text"] for b in data["blocks"]}
    assert "活动方案" in texts.values()
    # 找到"周末满 100 送 30。"那块的 id
    target = next(b["id"] for b in data["blocks"] if "周末满" in b["text"])
    canvas_docedit.write_blocks(p, {target: "周末满 200 送 80，力度加倍！"})

    after = {b["text"] for b in canvas_docedit.read_blocks(p)["blocks"]}
    assert "周末满 200 送 80，力度加倍！" in after
    assert "周末满 100 送 30。" not in after
    assert "仅限堂食。" in after      # 没动的块原样
    assert "活动方案" in after        # 标题原样
    # 改前自动备份了（落在内容库 .backups）
    assert (tmp_path / "lib" / ".backups").exists()


def test_pptx_read_edit_writeback(tmp_path):
    p = tmp_path / "培训.pptx"
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.shapes.title.text = "助教服务标准"
    body = slide.placeholders[1].text_frame
    body.text = "热情主动"
    body.add_paragraph().text = "不越界"
    prs.save(str(p))

    data = canvas_docedit.read_blocks(p)
    assert data["kind"] == "pptx"
    target = next(b["id"] for b in data["blocks"] if "热情主动" in b["text"])
    canvas_docedit.write_blocks(p, {target: "热情主动、有分寸"})

    after = {b["text"] for b in canvas_docedit.read_blocks(p)["blocks"]}
    assert "热情主动、有分寸" in after
    assert "热情主动" not in after
    assert "不越界" in after            # 没动的块原样
    assert "助教服务标准" in after      # 标题原样


def test_writeback_empty_noop(tmp_path):
    p = tmp_path / "x.docx"
    d = docx.Document()
    d.add_paragraph("一")
    d.save(str(p))
    canvas_docedit.write_blocks(p, {})  # 不抛错
    assert "一" in {b["text"] for b in canvas_docedit.read_blocks(p)["blocks"]}


# ────────── #M6-3: 段落编辑保留内联格式（加粗/颜色/字号） ──────────

def test_docx_edit_preserves_bold_run(tmp_path, monkeypatch):
    """部分编辑段落文字时，未改动 run 的加粗格式应保留。"""
    monkeypatch.setenv("DESKTOP_LIBRARY_DIR", str(tmp_path / "lib"))
    p = tmp_path / "bold.docx"
    d = docx.Document()
    para = d.add_paragraph()
    run1 = para.add_run("半价畅打 ")
    run2 = para.add_run("超值体验")
    run2.bold = True
    d.save(str(p))

    data = canvas_docedit.read_blocks(p)
    target = next(b for b in data["blocks"] if "半价畅打" in b["text"])
    canvas_docedit.write_blocks(p, {target["id"]: "半价畅打 极致体验"})

    d2 = docx.Document(str(p))
    para2 = [pp for pp in d2.paragraphs if pp.text.strip()][0]
    assert para2.text == "半价畅打 极致体验"
    runs = para2.runs
    assert len(runs) >= 2
    assert runs[0].bold is not True
    assert runs[1].bold is True


def test_docx_edit_preserves_font_size(tmp_path, monkeypatch):
    """编辑只改后半段时，前段的字号应原封不动。"""
    from docx.shared import Pt
    monkeypatch.setenv("DESKTOP_LIBRARY_DIR", str(tmp_path / "lib"))
    p = tmp_path / "size.docx"
    d = docx.Document()
    para = d.add_paragraph()
    r1 = para.add_run("标题")
    r1.font.size = Pt(24)
    r2 = para.add_run("正文内容")
    r2.font.size = Pt(12)
    d.save(str(p))

    data = canvas_docedit.read_blocks(p)
    target = next(b for b in data["blocks"] if "标题" in b["text"])
    canvas_docedit.write_blocks(p, {target["id"]: "标题改过的正文"})

    d2 = docx.Document(str(p))
    para2 = [pp for pp in d2.paragraphs if pp.text.strip()][0]
    assert para2.text == "标题改过的正文"
    assert para2.runs[0].font.size == Pt(24)
    assert para2.runs[0].text == "标题"


# ────────── #M6-4: 段落编辑保留超链接 ──────────

def test_docx_edit_preserves_hyperlink(tmp_path, monkeypatch):
    """编辑含超链接段落的非链接部分时，超链接元素和文字应保留。"""
    from docx.oxml.ns import qn
    from docx.oxml import OxmlElement
    monkeypatch.setenv("DESKTOP_LIBRARY_DIR", str(tmp_path / "lib"))

    p = tmp_path / "link.docx"
    d = docx.Document()
    para = d.add_paragraph()
    para.add_run("点击 ")
    # 手动构建超链接（python-docx 无高层 API）
    hyperlink = OxmlElement("w:hyperlink")
    hyperlink.set(qn("r:id"), "rId99")
    run_el = OxmlElement("w:r")
    t_el = OxmlElement("w:t")
    t_el.text = "这里"
    run_el.append(t_el)
    hyperlink.append(run_el)
    para._element.append(hyperlink)
    para.add_run(" 查看详情")
    d.save(str(p))

    data = canvas_docedit.read_blocks(p)
    target = next(b for b in data["blocks"] if "点击" in b["text"])
    canvas_docedit.write_blocks(p, {target["id"]: "点击 这里 了解更多"})

    d2 = docx.Document(str(p))
    para2 = [pp for pp in d2.paragraphs if "点击" in (pp.text or "")][0]
    assert "了解更多" in para2.text
    # 超链接元素仍在、文字未变
    hlinks = para2._element.findall(qn("w:hyperlink"))
    assert len(hlinks) >= 1
    hlink_runs = hlinks[0].findall(qn("w:r"))
    hlink_text = "".join((r.find(qn("w:t")).text or "") for r in hlink_runs if r.find(qn("w:t")) is not None)
    assert hlink_text == "这里"
