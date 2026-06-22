"""画板文档预览（/canvas/doc）的纯函数单测：Word→HTML 片段、PPT→逐页大纲。

只测不依赖 DB/鉴权的转换逻辑（端点本身需 DESKTOP_LOCAL + 鉴权，归 e2e）。"""
import docx
import pytest
from pptx import Presentation
from pptx.util import Inches

from api.v1.canvas import _docx_to_html, _pptx_to_slides, _safe_doc
from fastapi import HTTPException


def test_docx_to_html_headings_paragraphs_table(tmp_path):
    p = tmp_path / "方案.docx"
    d = docx.Document()
    d.add_heading("活动方案", level=1)
    d.add_paragraph("周末满 100 送 30。")
    d.add_heading("细则", level=2)
    d.add_paragraph("仅限堂食。")
    t = d.add_table(rows=1, cols=2)
    t.rows[0].cells[0].text = "项目"
    t.rows[0].cells[1].text = "金额"
    d.save(str(p))

    html = _docx_to_html(p)
    assert "<h1>活动方案</h1>" in html
    assert "<h2>细则</h2>" in html
    assert "周末满 100 送 30。" in html
    assert "<table" in html and "项目" in html and "金额" in html


def test_docx_to_html_escapes_html(tmp_path):
    p = tmp_path / "x.docx"
    d = docx.Document()
    d.add_paragraph("<script>alert(1)</script> & 来电")
    d.save(str(p))
    html = _docx_to_html(p)
    # 必须转义，绝不能把用户文档里的 <script> 原样吐成可执行标签
    assert "<script>" not in html
    assert "&lt;script&gt;" in html
    assert "&amp; 来电" in html


def test_pptx_to_slides_title_and_bullets(tmp_path):
    p = tmp_path / "培训.pptx"
    prs = Presentation()
    layout = prs.slide_layouts[1]  # 标题 + 内容
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = "助教服务标准"
    body = slide.placeholders[1].text_frame
    body.text = "热情主动"
    body.add_paragraph().text = "不越界"
    prs.save(str(p))

    slides = _pptx_to_slides(p)
    assert len(slides) == 1
    assert slides[0]["title"] == "助教服务标准"
    assert "热情主动" in slides[0]["bullets"]
    assert "不越界" in slides[0]["bullets"]
    # 标题不应混进要点里（标题 shape 被排除）
    assert "助教服务标准" not in slides[0]["bullets"]


def test_safe_doc_rejects_non_doc_suffix(tmp_path):
    f = tmp_path / "a.txt"
    f.write_text("hi")
    with pytest.raises(HTTPException) as ei:
        _safe_doc(str(f), [str(f)])
    assert ei.value.status_code == 400


def test_safe_doc_rejects_outside_sandbox(tmp_path):
    # 没在 selected_files 白名单里、也不在内容库 → 越界 403
    f = tmp_path / "secret.pdf"
    f.write_bytes(b"%PDF-1.4 fake")
    with pytest.raises(HTTPException) as ei:
        _safe_doc(str(f), [])
    assert ei.value.status_code == 403
