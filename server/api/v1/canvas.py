"""Canvas（画布）路由——成品右侧展开"指着某处说改这里"，报表可视化看/改，以及文档(PDF/Word/PPT/网页)右侧预览。"""
import base64
import os
from html import escape as _esc
from pathlib import Path
from types import SimpleNamespace
from typing import Annotated

from fastapi import APIRouter, Depends, HTTPException
from pydantic import BaseModel
from sqlalchemy.ext.asyncio import AsyncSession

from api.deps import get_current_user, get_current_store, get_db
from core.rbac import Permission, require_permission
from models.store import Store
from models.user import User
from services.agent.local_tools import _resolve as _sandbox_resolve
from services.canvas_service import canvas_edit

router = APIRouter()


def _require_desktop() -> None:
    """报表看/改读写本机文件——只在桌面本地版放行；云端多租户绝不能读服务器/他人文件。"""
    if os.environ.get("DESKTOP_LOCAL") != "1":
        raise HTTPException(status_code=403, detail="该功能仅桌面本地版可用")


def _safe_xlsx(path_str: str, selected_files: list[str] | None = None, full_disk: bool = False) -> Path:
    """校验并收敛报表路径：只允许操作【老板当场选定(OS 选择器授权)的文件】或内容库内，
    复用 local_tools 同款沙箱——防路径被下游替换成任意 xlsx(如 ~/账目.xlsx)。"""
    p = Path(path_str or "")
    if p.suffix.lower() not in (".xlsx", ".xlsm"):
        raise HTTPException(status_code=400, detail="只支持 .xlsx 报表")
    ctx = SimpleNamespace(allowed_paths=selected_files or [], full_disk_access=full_disk)
    try:
        resolved = _sandbox_resolve(path_str, ctx)
    except ValueError:
        raise HTTPException(status_code=403, detail="只能看/改你当场选定的报表")
    if not resolved.exists() or not resolved.is_file():
        raise HTTPException(status_code=404, detail="报表文件不存在")
    return resolved


class CanvasEditRequest(BaseModel):
    content: str                          # 当前成品全文（前端持有的最新版）
    instruction: str                      # 怎么改（老板的话）
    selection: str | None = None          # 老板圈中要改的那一段；空=整篇修订
    deliverable_type: str | None = None   # 成品类型(文案/活动方案/话术…)，只影响提示语气


class RenderRequest(BaseModel):
    content: str          # 定稿内容（Markdown/文本）
    format: str = "md"    # 导出格式：md/txt/html/docx（见 canvas_io.SUPPORTED）


@router.post("/render")
async def render_deliverable(
    body: RenderRequest,
    _user: Annotated[User, Depends(get_current_user)],
    _store: Annotated[Store, Depends(get_current_store)],
    _perm: None = Depends(require_permission(Permission.GENERATION_CREATE)),
):
    """把成品渲染成指定格式的字节(base64)，给前端走系统「另存为」写到本机任意位置。"""
    from services.canvas_io import SUPPORTED, render_bytes
    if body.format not in SUPPORTED:
        raise HTTPException(status_code=400, detail=f"不支持的格式：{body.format}")
    try:
        data = render_bytes(body.content, body.format)
    except Exception as e:  # noqa: BLE001
        raise HTTPException(status_code=422, detail=f"导出失败：{e}")
    return {"base64": base64.b64encode(data).decode("ascii"), "ext": body.format}


class SaveLibraryRequest(BaseModel):
    content: str
    format: str = "md"
    name: str = "成品"     # 文件名（不含扩展名）


@router.post("/save-to-library")
async def save_to_library_endpoint(
    body: SaveLibraryRequest,
    _user: Annotated[User, Depends(get_current_user)],
    _store: Annotated[Store, Depends(get_current_store)],
    _perm: None = Depends(require_permission(Permission.GENERATION_CREATE)),
):
    """把成品存进「内容库/成品」（桌面专属，重名自动备份）。返回写入路径。"""
    _require_desktop()
    from services.canvas_io import SUPPORTED, save_to_library
    if body.format not in SUPPORTED:
        raise HTTPException(status_code=400, detail=f"不支持的格式：{body.format}")
    try:
        path = save_to_library(body.content, body.format, body.name)
    except Exception as e:  # noqa: BLE001
        raise HTTPException(status_code=422, detail=f"保存失败：{e}")
    return {"ok": True, "path": str(path)}


@router.post("/edit")
async def canvas_edit_endpoint(
    body: CanvasEditRequest,
    user: Annotated[User, Depends(get_current_user)],
    store: Annotated[Store, Depends(get_current_store)],
    db: Annotated[AsyncSession, Depends(get_db)],
    _perm: None = Depends(require_permission(Permission.GENERATION_CREATE)),
):
    """圈了段只改那段（改这里不动别处），没圈则整篇修订。复用统一管道(配额/合规/落库/BYOK)。"""
    return await canvas_edit(
        db, store, user,
        content=body.content,
        instruction=body.instruction,
        selection=body.selection,
        deliverable_type=body.deliverable_type or "内容",
    )


# ───────────────── 报表可视化：看表格（只读）/ 点格改（写·自动备份） ·桌面专属 ─────────────────

_MAX_ROWS = 200
_MAX_COLS = 30


class SheetRequest(BaseModel):
    path: str  # 老板选定的本机报表路径(.xlsx)
    selected_files: list[str] = []   # OS 选择器授权的文件(沙箱白名单)；path 必须在其中或内容库内
    full_disk_access: bool = False


@router.post("/sheet")
async def read_sheet(
    body: SheetRequest,
    _user: Annotated[User, Depends(get_current_user)],
    _store: Annotated[Store, Depends(get_current_store)],
    _perm: None = Depends(require_permission(Permission.QUOTA_VIEW)),
):
    """读本机报表 → 结构化表格(供前端可视化展示)。桌面专属、只读、限行列防超大。"""
    _require_desktop()
    p = _safe_xlsx(body.path, body.selected_files, body.full_disk_access)
    from openpyxl import load_workbook
    wb = load_workbook(p, data_only=True)
    truncated = len(wb.worksheets) > 5
    sheets = []
    for ws in wb.worksheets[:5]:
        if (ws.max_row or 0) > _MAX_ROWS or (ws.max_column or 0) > _MAX_COLS:
            truncated = True
        rows = []
        for row in ws.iter_rows(max_row=_MAX_ROWS, max_col=_MAX_COLS):
            rows.append(["" if c.value is None else str(c.value) for c in row])
        while rows and not any(cell for cell in rows[-1]):  # 去尾部全空行
            rows.pop()
        # 去尾部全空列（openpyxl 按 max_col 填充会带一堆空列）
        maxc = 0
        for row in rows:
            for i in range(len(row) - 1, -1, -1):
                if row[i]:
                    maxc = max(maxc, i + 1)
                    break
        rows = [row[:maxc] for row in rows]
        sheets.append({"name": ws.title, "rows": rows})
    wb.close()
    return {"name": p.name, "sheets": sheets, "truncated": truncated}


class ExcelCellEdit(BaseModel):
    path: str
    cell: str                 # A1 式坐标，如 B2
    value: str                # 新值（纯数字会自动转数值）
    sheet: str | None = None  # 工作表名，留空=第一个表
    selected_files: list[str] = []   # OS 选择器授权的文件(沙箱白名单)
    full_disk_access: bool = False


@router.post("/excel-edit")
async def excel_edit(
    body: ExcelCellEdit,
    _user: Annotated[User, Depends(get_current_user)],
    _store: Annotated[Store, Depends(get_current_store)],
    _perm: None = Depends(require_permission(Permission.GENERATION_CREATE)),
):
    """点格改：改本机报表一个单元格。桌面专属、改前自动备份、返回前后对比(diff)。"""
    _require_desktop()
    import shutil
    from core.timezone import business_now
    from openpyxl import load_workbook
    p = _safe_xlsx(body.path, body.selected_files, body.full_disk_access)
    # 改前自动备份（落到报表同目录的隐藏备份夹，可回滚）
    bdir = p.parent / ".billiards-backups"
    bdir.mkdir(exist_ok=True)
    backup = bdir / f"{p.stem}.{business_now().strftime('%Y%m%d-%H%M%S')}{p.suffix}.bak"
    wb = load_workbook(p)
    ws = wb[body.sheet] if (body.sheet and body.sheet in wb.sheetnames) else wb.active
    try:
        old = ws[body.cell].value
    except Exception:
        wb.close()
        raise HTTPException(status_code=400, detail=f"坐标无效：{body.cell}")
    # 纯数字转数值，否则按文本
    v: object = body.value
    s = (body.value or "").strip()
    try:
        v = int(s) if s.lstrip("-").isdigit() else float(s)
    except ValueError:
        v = body.value
    try:
        ws[body.cell] = v
    except AttributeError:  # 合并单元格的非左上角格只读
        wb.close()
        raise HTTPException(status_code=400, detail=f"{body.cell} 是合并单元格，请改合并区左上角那个格子")
    shutil.copy2(p, backup)  # 改成功才备份，避免改失败留孤儿备份
    wb.save(p)
    wb.close()
    return {"ok": True, "sheet": ws.title, "cell": body.cell,
            "old": "" if old is None else str(old), "new": str(v)}


# ───────────────── 文档预览：PDF / Word(.docx) / PPT(.pptx) / 网页(.html) ·桌面专属·只读 ─────────────────
# 学 Claude Artifacts / ChatGPT Canvas：成品/模板在右侧原样看清楚，老板确认后再定稿。
# 沿用 readSheet 的套路——后端读本地文件 → 返回前端可渲染的数据（前端不能直接读任意本地文件字节）。

_DOC_SUFFIXES = (".pdf", ".docx", ".pptx", ".html", ".htm")
_DOC_MAX_PDF_BYTES = 20 * 1024 * 1024   # PDF 走 base64 内联，>20MB 不预览（让老板用系统阅读器开）
_DOC_MAX_HTML_CHARS = 2_000_000         # 网页原文上限，防超大 srcDoc
_DOC_MAX_SLIDES = 80                     # PPT 最多预览页数
_DOC_MAX_BULLETS = 60                    # 每页最多要点数


def _safe_doc(path_str: str, selected_files: list[str] | None = None, full_disk: bool = False) -> Path:
    """校验并收敛文档路径（复用 local_tools 沙箱），只允许可预览的文档类型。"""
    p = Path(path_str or "")
    if p.suffix.lower() not in _DOC_SUFFIXES:
        raise HTTPException(status_code=400, detail="只支持预览 PDF / Word(.docx) / PPT(.pptx) / 网页(.html)")
    ctx = SimpleNamespace(allowed_paths=selected_files or [], full_disk_access=full_disk)
    try:
        resolved = _sandbox_resolve(path_str, ctx)
    except ValueError:
        raise HTTPException(status_code=403, detail="只能预览你当场选定的文件")
    if not resolved.exists() or not resolved.is_file():
        raise HTTPException(status_code=404, detail="文件不存在")
    return resolved


def _docx_to_html(p: Path) -> str:
    """Word(.docx) → 安全的 HTML 片段（标题/段落/加粗斜体/列表/表格）。无第三方转换依赖，纯 python-docx。"""
    import docx
    doc = docx.Document(str(p))
    out: list[str] = []
    for para in doc.paragraphs:
        style = (para.style.name or "").lower() if para.style else ""
        runs_html = ""
        for run in para.runs:
            t = _esc(run.text)
            if not t:
                continue
            if run.bold:
                t = f"<strong>{t}</strong>"
            if run.italic:
                t = f"<em>{t}</em>"
            runs_html += t
        if not runs_html:
            if not para.text.strip():
                out.append("<p><br/></p>")
                continue
            runs_html = _esc(para.text)
        if style.startswith("heading 1") or style == "title":
            out.append(f"<h1>{runs_html}</h1>")
        elif style.startswith("heading 2"):
            out.append(f"<h2>{runs_html}</h2>")
        elif style.startswith("heading 3"):
            out.append(f"<h3>{runs_html}</h3>")
        elif "list" in style or style.startswith("bullet"):
            out.append(f"<li>{runs_html}</li>")
        else:
            out.append(f"<p>{runs_html}</p>")
    for tbl in doc.tables:
        rows_html = ""
        for row in tbl.rows:
            cells = "".join(f"<td>{_esc(c.text.strip())}</td>" for c in row.cells)
            if cells:
                rows_html += f"<tr>{cells}</tr>"
        if rows_html:
            out.append(f"<table border='1' cellspacing='0' cellpadding='6'>{rows_html}</table>")
    return "\n".join(out) if out else "<p>（这份 Word 文档没有文字内容）</p>"


def _pptx_to_slides(p: Path) -> list[dict]:
    """PPT(.pptx) → 逐页大纲（标题 + 要点）。轻量文字版，零重依赖；原样幻灯片渲染留作后续升级。"""
    from pptx import Presentation
    prs = Presentation(str(p))
    slides: list[dict] = []
    for slide in prs.slides:
        if len(slides) >= _DOC_MAX_SLIDES:
            break
        title = ""
        title_id = None  # python-pptx 每次访问 shape 都是新包装对象，用稳定的 shape_id 比对、别用 is
        try:
            if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
                title_id = slide.shapes.title.shape_id
                title = (slide.shapes.title.text_frame.text or "").strip()
        except Exception:  # noqa: BLE001
            title = ""
        bullets: list[str] = []
        for shape in slide.shapes:
            if not shape.has_text_frame or shape.shape_id == title_id:
                continue
            for para in shape.text_frame.paragraphs:
                t = ("".join(r.text for r in para.runs).strip() or (para.text or "").strip())
                if t:
                    bullets.append(t)
                    if len(bullets) >= _DOC_MAX_BULLETS:
                        break
            if len(bullets) >= _DOC_MAX_BULLETS:
                break
        slides.append({"title": title, "bullets": bullets})
    return slides


class DocRequest(BaseModel):
    path: str
    selected_files: list[str] = []   # OS 选择器授权的文件(沙箱白名单)
    full_disk_access: bool = False


@router.post("/doc")
async def read_doc(
    body: DocRequest,
    _user: Annotated[User, Depends(get_current_user)],
    _store: Annotated[Store, Depends(get_current_store)],
    _perm: None = Depends(require_permission(Permission.QUOTA_VIEW)),
):
    """读本机文档 → 前端可渲染的数据。桌面专属、只读。
    返回 render：pdf(base64原样) / page(网页原文) / richtext(Word转HTML片段) / slides(PPT逐页大纲) / toobig。"""
    _require_desktop()
    p = _safe_doc(body.path, body.selected_files, body.full_disk_access)
    ext = p.suffix.lower()
    try:
        if ext == ".pdf":
            data = p.read_bytes()
            if len(data) > _DOC_MAX_PDF_BYTES:
                return {"name": p.name, "render": "toobig",
                        "message": f"这份 PDF 约 {len(data) // 1024 // 1024}MB，超过预览上限，请用系统自带阅读器打开"}
            return {"name": p.name, "render": "pdf", "pdf_base64": base64.b64encode(data).decode("ascii")}
        if ext in (".html", ".htm"):
            text = p.read_text(encoding="utf-8", errors="replace")
            truncated = len(text) > _DOC_MAX_HTML_CHARS
            return {"name": p.name, "render": "page", "html": text[:_DOC_MAX_HTML_CHARS], "truncated": truncated}
        if ext == ".docx":
            return {"name": p.name, "render": "richtext", "html": _docx_to_html(p)}
        if ext == ".pptx":
            slides = _pptx_to_slides(p)
            return {"name": p.name, "render": "slides", "slides": slides, "truncated": len(slides) >= _DOC_MAX_SLIDES}
    except HTTPException:
        raise
    except Exception as e:  # noqa: BLE001
        raise HTTPException(status_code=422, detail=f"这个文档读不开（可能损坏或格式特殊）：{e}")
    raise HTTPException(status_code=400, detail="不支持预览这种文档")


# ───────────────── Word/PPT 文字级编辑 + 写回原文件 ·桌面专属 ─────────────────

def _safe_editable_doc(path_str: str, selected_files: list[str] | None, full_disk: bool) -> Path:
    """收敛路径并限定为可按块编辑的格式（docx/pptx）。"""
    p = _safe_doc(path_str, selected_files, full_disk)
    if p.suffix.lower() not in (".docx", ".pptx"):
        raise HTTPException(status_code=400, detail="只支持 Word(.docx) / PPT(.pptx) 按块编辑")
    return p


@router.post("/doc-blocks")
async def read_doc_blocks(
    body: DocRequest,
    _user: Annotated[User, Depends(get_current_user)],
    _store: Annotated[Store, Depends(get_current_store)],
    _perm: None = Depends(require_permission(Permission.QUOTA_VIEW)),
):
    """读 Word/PPT → 带稳定 id 的文本块（供前端逐块编辑）。"""
    _require_desktop()
    p = _safe_editable_doc(body.path, body.selected_files, body.full_disk_access)
    from services.canvas_docedit import read_blocks
    try:
        data = read_blocks(p)
    except Exception as e:  # noqa: BLE001
        raise HTTPException(status_code=422, detail=f"读不开这个文档：{e}")
    return {"name": p.name, **data}


class DocSaveRequest(BaseModel):
    path: str
    edits: dict[str, str]            # {块id: 新文字}
    selected_files: list[str] = []
    full_disk_access: bool = False


@router.post("/doc-save")
async def save_doc_blocks(
    body: DocSaveRequest,
    _user: Annotated[User, Depends(get_current_user)],
    _store: Annotated[Store, Depends(get_current_store)],
    _perm: None = Depends(require_permission(Permission.GENERATION_CREATE)),
):
    """把改动按 id 原地写回原 Word/PPT（保留结构、改前自动备份）。"""
    _require_desktop()
    p = _safe_editable_doc(body.path, body.selected_files, body.full_disk_access)
    from services.canvas_docedit import write_blocks
    try:
        write_blocks(p, body.edits or {})
    except Exception as e:  # noqa: BLE001
        raise HTTPException(status_code=422, detail=f"写回失败：{e}")
    return {"ok": True, "path": str(p), "saved": len(body.edits or {})}
